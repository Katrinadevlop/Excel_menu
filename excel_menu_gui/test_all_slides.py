#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Тест всех слайдов презентации
"""
import os
import sys
from pathlib import Path

sys.path.append(os.path.dirname(os.path.abspath(__file__)))

from presentation_handler import create_presentation_with_fish_and_side_dishes

def analyze_all_slides(pptx_path: str):
    """Анализирует все слайды созданной презентации"""
    print(f"\n🔍 АНАЛИЗ ВСЕХ СЛАЙДОВ ПРЕЗЕНТАЦИИ")
    print("=" * 70)
    
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        prs = Presentation(pptx_path)
        print(f"📊 Общее количество слайдов: {len(prs.slides)}")
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            print(f"\n🎯 СЛАЙД {slide_idx}:")
            
            tables_found = 0
            all_text = []
            
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    tables_found += 1
                    table = shape.table
                    print(f"\n📋 Таблица {tables_found}: {len(table.rows)} строк × {len(table.columns)} столбцов")
                    
                    # Показываем первые несколько строк таблицы
                    max_rows_to_show = min(5, len(table.rows))
                    for i in range(max_rows_to_show):
                        row = table.rows[i]
                        row_content = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            all_text.append(cell_text)
                            if cell_text:
                                row_content.append(f"'{cell_text}'")
                            else:
                                row_content.append("'[пусто]'")
                        print(f"   Строка {i+1}: {' | '.join(row_content)}")
                    
                    if len(table.rows) > max_rows_to_show:
                        print(f"   ... (еще {len(table.rows) - max_rows_to_show} строк)")
                elif hasattr(shape, 'text'):
                    # Собираем весь текст со слайда
                    text = shape.text.strip()
                    if text:
                        all_text.append(text)
            
            if not tables_found:
                print("   📋 Таблиц не найдено")
                # Показываем весь текст если нет таблиц
                slide_text = ' '.join(all_text)
                if slide_text:
                    print(f"   📝 Текст на слайде: {slide_text[:200]}...")
            
            # Проверяем на наличие блюд (ключевые слова)
            slide_text = ' '.join(all_text).upper()
            dish_indicators = ['КОТЛЕТА', 'СУП', 'БОРЩ', 'КАША', 'РЫБНАЯ', 'ЖАРЕНАЯ', 'САЛАТ', 'ЗАКУСКА']
            found_dishes = [word for word in dish_indicators if word in slide_text]
            
            if found_dishes:
                print(f"   🍽️  Найденные блюда: {', '.join(found_dishes[:3])}{'...' if len(found_dishes) > 3 else ''}")
            else:
                print("   ❌ Блюда не обнаружены")
                
    except ImportError:
        print("❌ Модуль python-pptx не установлен")
    except Exception as e:
        print(f"❌ Ошибка анализа презентации: {e}")
        import traceback
        traceback.print_exc()

def test_all_slides():
    """Тестирует создание презентации и анализирует все слайды"""
    print("🧪 ТЕСТ ВСЕХ СЛАЙДОВ ПРЕЗЕНТАЦИИ")
    print("=" * 70)
    
    # Используем тот же проблемный файл
    excel_path = Path(r"C:\Users\katya\Downloads\Telegram Desktop\4 сентября - четверг (2).xls")
    template_path = Path("templates/presentation_template.pptx")
    output_path = Path("test_all_slides_presentation.pptx")
    
    if not excel_path.exists():
        print(f"❌ Excel файл не найден: {excel_path}")
        return
        
    if not template_path.exists():
        print(f"❌ Шаблон не найден: {template_path}")
        return
    
    print(f"📄 Excel файл: {excel_path.name}")
    print(f"📄 Шаблон: {template_path}")
    print(f"💾 Результат: {output_path}")
    
    # Создаем презентацию
    print(f"\n🎯 СОЗДАНИЕ ПРЕЗЕНТАЦИИ")
    try:
        success, message = create_presentation_with_fish_and_side_dishes(
            str(template_path),
            str(excel_path),
            str(output_path)
        )
        
        if success:
            print(f"✅ Презентация создана: {output_path}")
            print(f"📝 Сообщение: {message}")
            
            if output_path.exists():
                size = output_path.stat().st_size
                print(f"📏 Размер файла: {size:,} байт")
                
                # Анализируем все слайды
                analyze_all_slides(str(output_path))
        else:
            print(f"❌ Ошибка создания презентации: {message}")
            
    except Exception as e:
        print(f"❌ Исключение: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    test_all_slides()
