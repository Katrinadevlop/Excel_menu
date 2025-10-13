import re
import shutil
from pathlib import Path
from typing import List, Tuple, Optional
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.dish_extractor import (
    DishItem,
    extract_salads_from_excel,
    extract_first_courses_from_excel,
    extract_meat_dishes_from_excel,
    extract_poultry_dishes_from_excel,
    extract_fish_dishes_from_column_e,
    extract_fish_dishes_from_excel,
    extract_side_dishes_from_excel, extract_dishes_from_excel,
)


def update_slide_with_dishes(slide, dishes: List[DishItem]) -> bool:
    """
    Обновляет один слайд презентации, вставляя данные блюд в подходящую таблицу и форматируя содержимое.

    Args:
        slide (pptx.slide.Slide): Слайд, в котором требуется обновить таблицу.
        dishes (List[DishItem]): Список блюд для вставки (name, weight, price).

    Returns:
        bool: True, если таблица найдена и успешно обновлена; иначе False.
    """
    try:
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        # Найдем все таблицы на слайде
        table_shapes = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_shapes.append(shape)
                
        if not table_shapes:
            print("На слайде не найдено таблиц")
            return False
            
        # Умный выбор таблицы: предпочитаем таблицы с правильными заголовками
        best_table_shape = None
        best_score = -1
        
        for shape in table_shapes:
            table = shape.table
            rows = len(table.rows)
            data_rows = rows - 1 if rows > 1 else 0
            
            if data_rows <= 0:
                continue
                
            score = 0
            
            # Анализируем заголовки первой строки
            if rows > 0:
                header_row = table.rows[0]
                headers = [cell.text.strip() for cell in header_row.cells]
                
                # Предпочитаем таблицы с пустым первым заголовком (для названий блюд)
                if len(headers) >= 3:
                    first_header = headers[0].upper()
                    second_header = headers[1].upper()
                    third_header = headers[2].upper()
                    
                    # Хороший заголовок: пустое поле для названия + вес + цена
                    if (not first_header or first_header == '') and \
                       ('ВЕС' in second_header or 'ГРАММ' in second_header) and \
                       ('ЦЕНА' in third_header or 'РУБ' in third_header):
                        score += 100  # Отличная таблица
                    # Плохой заголовок: уже заполнен данными
                    elif len(first_header) > 5 and any(char.isalpha() for char in first_header):
                        score -= 50  # Плохая таблица (уже с данными)
            
            # Добавляем баллы за количество строк
            score += data_rows
            
            
            if score > best_score:
                best_score = score
                best_table_shape = shape
                max_data_rows = data_rows
        
        if best_table_shape is None:
            return False
            
        
        table = best_table_shape.table
        
        # Получаем количество строк в таблице
        total_rows = len(table.rows)
        
        # Определяем оптимальный размер шрифта (начинаем с 28pt)
        available_rows = total_rows - 1  # -1 для заголовка
        dishes_to_show = len(dishes)
        
        # Начинаем с 28pt и уменьшаем, если нужно
        if dishes_to_show <= available_rows:
            font_size = 28  # Оптимальный размер
        elif dishes_to_show <= available_rows * 1.5:
            font_size = 24  # Немного уменьшаем
        elif dishes_to_show <= available_rows * 2:
            font_size = 20  # Еще уменьшаем
        elif dishes_to_show <= available_rows * 3:
            font_size = 16  # Минимально читаемый
        else:
            font_size = 14  # Критически маленький
        
        # Ограничиваем количество блюд доступным местом
        dishes_to_fill = dishes[:available_rows]
        
        # Очищаем все строки кроме первой (заголовки) и заполняем их блюдами
        for i, dish in enumerate(dishes_to_fill):
            row_idx = i + 1  # +1 потому что 0 - это заголовок
            
            if row_idx < total_rows:
                row = table.rows[row_idx]
                
                # Заполняем и форматируем ячейки
                if len(row.cells) >= 3:
                    # Название блюда
                    cell_name = row.cells[0]
                    cell_name.text = dish.name
                    if cell_name.text_frame.paragraphs:
                        paragraph = cell_name.text_frame.paragraphs[0]
                        paragraph.alignment = PP_ALIGN.LEFT
                        # Устанавливаем отступ в 10 пикселей
                        cell_name.text_frame.margin_left = Pt(10)
                        cell_name.text_frame.margin_right = Pt(10)
                        cell_name.text_frame.margin_top = Pt(10)
                        cell_name.text_frame.margin_bottom = Pt(10)
                        if paragraph.runs:
                            run = paragraph.runs[0]
                            run.font.name = 'Gilroy Medium'
                            run.font.size = Pt(font_size)
                            run.font.color.rgb = RGBColor(255, 255, 255)  # Белый цвет
                    
                    # Вес/объем
                    cell_weight = row.cells[1]
                    cell_weight.text = dish.weight
                    if cell_weight.text_frame.paragraphs:
                        paragraph = cell_weight.text_frame.paragraphs[0]
                        paragraph.alignment = PP_ALIGN.CENTER
                        # Устанавливаем отступ в 10 пикселей
                        cell_weight.text_frame.margin_left = Pt(10)
                        cell_weight.text_frame.margin_right = Pt(10)
                        cell_weight.text_frame.margin_top = Pt(10)
                        cell_weight.text_frame.margin_bottom = Pt(10)
                        if paragraph.runs:
                            run = paragraph.runs[0]
                            run.font.name = 'Gilroy Medium'
                            run.font.size = Pt(font_size)
                            run.font.color.rgb = RGBColor(255, 255, 255)
                    
                    # Цена
                    cell_price = row.cells[2]
                    # Убираем "руб." и другие обозначения валюты из цены для презентации
                    price_text = dish.price
                    # Убираем различные варианты написания рублей
                    price_text = re.sub(r'\s*(руб\.?|рублей|р\.?|₽|RUB)', '', price_text, flags=re.IGNORECASE)
                    price_text = price_text.strip()
                    cell_price.text = price_text
                    if cell_price.text_frame.paragraphs:
                        paragraph = cell_price.text_frame.paragraphs[0]
                        paragraph.alignment = PP_ALIGN.CENTER
                        # Устанавливаем отступ в 10 пикселей
                        cell_price.text_frame.margin_left = Pt(10)
                        cell_price.text_frame.margin_right = Pt(10)
                        cell_price.text_frame.margin_top = Pt(10)
                        cell_price.text_frame.margin_bottom = Pt(10)
                        if paragraph.runs:
                            run = paragraph.runs[0]
                            run.font.name = 'Gilroy Medium'
                            run.font.size = Pt(font_size)
                            run.font.color.rgb = RGBColor(255, 255, 255)
        
        # Очищаем оставшиеся строки
        for i in range(len(dishes_to_fill) + 1, total_rows):
            if i < len(table.rows):
                row = table.rows[i]
                for j in range(len(row.cells)):
                    row.cells[j].text = ""
                
        return True
        
    except Exception as e:
        return False


def update_presentation_with_all_categories(presentation_path: str, all_dishes: dict, output_path: str) -> bool:
    """
    Обновляет презентацию, вставляя данные всех категорий блюд в соответствующие слайды.

    Args:
        presentation_path (str): Путь к исходному шаблону презентации (.pptx).
        all_dishes (dict): Данные по категориям:
            {'salads': List[DishItem], 'first_courses': List[DishItem], 'meat': List[DishItem],
             'poultry': List[DishItem], 'fish': List[DishItem], 'side_dishes': List[DishItem]}.
        output_path (str): Путь для сохранения обновлённой презентации.

    Returns:
        bool: True, если хотя бы один слайд успешно обновлён; иначе False.
    """
    try:
        # Копируем исходную презентацию
        shutil.copy2(presentation_path, output_path)
        
        # Открываем презентацию
        prs = Presentation(output_path)
        
        # Проверяем, что у нас достаточно слайдов
        if len(prs.slides) < 7:
            return False
            
        # Маппинг категорий на индексы слайдов
        slide_mapping = {
            'salads': 1,        # Слайд 2: САЛАТЫ И ХОЛОДНЫЕ ЗАКУСКИ
            'first_courses': 2, # Слайд 3: ПЕРВЫЕ БЛЮДА
            'meat': 3,          # Слайд 4: МЯСНЫЕ БЛЮДА
            'poultry': 4,       # Слайд 5: БЛЮДА ИЗ ПТИЦЫ
            'fish': 5,          # Слайд 6: РЫБНЫЕ БЛЮДА
            'side_dishes': 6    # Слайд 7: ГАРНИРЫ
        }
        
        # Обновляем каждый слайд соответствующими данными
        success_count = 0
        for category, slide_idx in slide_mapping.items():
            if category in all_dishes and all_dishes[category]:
                slide = prs.slides[slide_idx]
                if update_slide_with_dishes(slide, all_dishes[category]):
                    success_count += 1
                
        # Сохраняем презентацию
        prs.save(output_path)
        
        return success_count > 0
        
    except Exception as e:
        return False


def create_presentation_with_fish_and_side_dishes(template_path: str, excel_path: str, output_path: str) -> Tuple[bool, str]:
    """
    Создаёт презентацию и заполняет слайды рыбными блюдами (6-й слайд) и гарнирами (7-й слайд),
    используя общий механизм обновления презентации.

    Args:
        template_path (str): Путь к шаблону презентации .pptx.
        excel_path (str): Путь к Excel-файлу с меню.
        output_path (str): Путь для сохранения готовой презентации.

    Returns:
        Tuple[bool, str]: Пара (успех, сообщение).
    """
    try:
        if not Path(template_path).exists():
            return False, f"Шаблон презентации не найден: {template_path}"
        if not Path(excel_path).exists():
            return False, f"Excel файл не найден: {excel_path}"

        # Рыба: сначала пробуем столбец E (как раньше), затем общий поиск
        fish_dishes = extract_fish_dishes_from_column_e(excel_path)
        if len(fish_dishes) == 0:
            fish_dishes = extract_fish_dishes_from_excel(excel_path)

        # Гарниры
        side_dishes = extract_side_dishes_from_excel(excel_path)

        if len(fish_dishes) == 0 and len(side_dishes) == 0:
            return False, "В Excel файле не найдены рыбные блюда/гарниры. Проверьте структуру файла и названия категорий."

        # Собираем только нужные категории и вызываем общий обновитель
        all_dishes = {
            'salads': [],
            'first_courses': [],
            'meat': [],
            'poultry': [],
            'fish': fish_dishes,
            'side_dishes': side_dishes,
        }

        ok = update_presentation_with_all_categories(template_path, all_dishes, output_path)
        if not ok:
            return False, "Ошибка при обновлении презентации (рыба/гарниры)"

        parts = []
        if len(fish_dishes) > 0:
            parts.append(f"6-й слайд (рыба): {len(fish_dishes)} блюд")
        if len(side_dishes) > 0:
            parts.append(f"7-й слайд (гарниры): {len(side_dishes)} блюд")
        message = "Презентация создана!\n" + "\n".join(parts)
        return True, message

    except Exception as e:
        return False, f"Ошибка: {str(e)}"


def create_presentation_with_excel_data(template_path: str, excel_path: str, output_path: str) -> Tuple[bool, str]:
    """
    Создаёт презентацию и заполняет слайды данными всех категорий блюд.

    Args:
        template_path (str): Путь к шаблону презентации .pptx.
        excel_path (str): Путь к Excel-файлу с меню.
        output_path (str): Путь для сохранения готовой презентации.

    Returns:
        Tuple[bool, str]: Пара (успех, подробное сообщение о вставленных данных).
    """
    try:
        if not Path(template_path).exists():
            return False, f"Шаблон презентации не найден: {template_path}"
        if not Path(excel_path).exists():
            return False, f"Excel файл не найден: {excel_path}"
        
        print(f"🔍 Ищем салаты в файле: {excel_path}")
        salads = extract_salads_from_excel(excel_path)
        print(f"Салаты: найдено {len(salads)} блюд")
        if len(salads) == 0:
            print("Альтернативный поиск салатов...")
            salads = extract_dishes_from_excel(excel_path, ['САЛАТЫ', 'ХОЛОДНЫЕ ЗАКУСКИ', 'САЛАТЫ И ХОЛОДНЫЕ ЗАКУСКИ'])
            print(f"Салаты (альтернативный поиск): найдено {len(salads)} блюд")
        
        print(f"🔍 Ищем первые блюда в файле: {excel_path}")
        first_courses = extract_first_courses_from_excel(excel_path)
        if len(first_courses) == 0:
            # Фолбек на универсальный извлекатель
            first_courses = extract_dishes_from_excel(excel_path, ['ПЕРВЫЕ БЛЮДА', 'ПЕРВЫЕ'])
        
        meat_dishes = extract_meat_dishes_from_excel(excel_path)
        if len(meat_dishes) == 0:
            meat_dishes = extract_dishes_from_excel(excel_path, ['БЛЮДА ИЗ МЯСА', 'МЯСНЫЕ БЛЮДА'])
        
        poultry_dishes = extract_poultry_dishes_from_excel(excel_path)
        if len(poultry_dishes) == 0:
            poultry_dishes = extract_dishes_from_excel(excel_path, ['БЛЮДА ИЗ ПТИЦЫ', 'БЛЮДА ИЗ КУРИЦЫ', 'КУРИНЫЕ БЛЮДА'])
        
        fish_dishes = extract_fish_dishes_from_column_e(excel_path)
        if len(fish_dishes) == 0:
            fish_dishes = extract_fish_dishes_from_excel(excel_path)
        
        side_dishes = extract_side_dishes_from_excel(excel_path)
        if len(side_dishes) == 0:
            side_dishes = extract_dishes_from_excel(excel_path, ['ГАРНИРЫ', 'ГАРНИР'])
        
        
        total_dishes = len(salads) + len(first_courses) + len(meat_dishes) + len(poultry_dishes) + len(fish_dishes) + len(side_dishes)
        if total_dishes == 0:
            return False, "В Excel файле не найдены блюда указанных категорий. Проверьте структуру файла и названия категорий."

        all_dishes = {
            'salads': salads,
            'first_courses': first_courses,
            'meat': meat_dishes,
            'poultry': poultry_dishes,
            'fish': fish_dishes,
            'side_dishes': side_dishes,
        }

        success = update_presentation_with_all_categories(template_path, all_dishes, output_path)
        if success:
            results = []
            if len(salads) > 0:
                results.append(f"Салаты и холодные закуски: {len(salads)} блюд")
            if len(first_courses) > 0:
                results.append(f"Первые блюда: {len(first_courses)} блюд")
            if len(meat_dishes) > 0:
                results.append(f"Блюда из мяса: {len(meat_dishes)} блюд")
            if len(poultry_dishes) > 0:
                results.append(f"Блюда из птицы: {len(poultry_dishes)} блюд")
            if len(fish_dishes) > 0:
                results.append(f"Блюда из рыбы: {len(fish_dishes)} блюд")
            if len(side_dishes) > 0:
                results.append(f"Гарниры: {len(side_dishes)} блюд")
            message = "Презентация создана!\n" + "\n".join(results)
            return True, message
        else:
            return False, "Ошибка при обновлении презентации"
            
    except Exception as e:
        return False, f"Ошибка: {str(e)}"
