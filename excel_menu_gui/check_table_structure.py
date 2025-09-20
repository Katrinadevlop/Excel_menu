#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Проверяем структуру таблицы в шаблоне презентации
"""
import sys
import os
from pathlib import Path

sys.path.append(os.path.dirname(os.path.abspath(__file__)))

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    
    template_path = Path("templates/presentation_template.pptx")
    
    if not template_path.exists():
        print(f"❌ Шаблон не найден: {template_path}")
        sys.exit(1)
    
    print(f"📄 Анализируем шаблон: {template_path}")
    
    prs = Presentation(str(template_path))
    print(f"📊 Всего слайдов: {len(prs.slides)}")
    
    if len(prs.slides) >= 6:
        slide_6 = prs.slides[5]  # 6-й слайд
        print(f"\n🎯 АНАЛИЗ 6-ГО СЛАЙДА (рыбные блюда):")
        
        tables_found = 0
        for i, shape in enumerate(slide_6.shapes):
            print(f"  Элемент {i+1}: тип {shape.shape_type}")
            
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                tables_found += 1
                table = shape.table
                print(f"    📋 ТАБЛИЦА {tables_found}:")
                print(f"       Строк: {len(table.rows)}")
                print(f"       Столбцов: {len(table.columns)}")
                print(f"       Строк для данных: {len(table.rows) - 1} (без заголовка)")
                
                # Показываем содержимое первой строки (заголовки)
                if len(table.rows) > 0:
                    header_row = table.rows[0]
                    headers = []
                    for cell in header_row.cells:
                        headers.append(f"'{cell.text.strip()}'")
                    print(f"       Заголовки: {' | '.join(headers)}")
                
                # Показываем пример пустых строк для данных
                print(f"       Пустые строки для заполнения:")
                for row_idx in range(1, min(6, len(table.rows))):  # Показываем первые 5 строк данных
                    row = table.rows[row_idx]
                    row_content = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        row_content.append(f"'{cell_text}'" if cell_text else "'[пусто]'")
                    print(f"         Строка {row_idx}: {' | '.join(row_content)}")
                
                print()
        
        if tables_found == 0:
            print("  ❌ На 6-м слайде не найдено таблиц!")
        else:
            print(f"  ✅ Найдено таблиц: {tables_found}")
    else:
        print("❌ В презентации меньше 6 слайдов")

except ImportError:
    print("❌ Модуль python-pptx не установлен")
except Exception as e:
    print(f"❌ Ошибка: {e}")
