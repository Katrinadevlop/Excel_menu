#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Тест создания презентации с проблемным файлом
"""
import os
import sys
from pathlib import Path

sys.path.append(os.path.dirname(os.path.abspath(__file__)))

from presentation_handler import (
    extract_fish_dishes_from_column_e,
    create_presentation_with_fish_and_side_dishes
)

def test_problematic_file_presentation():
    """Тестирует создание презентации с проблемным файлом"""
    print("🧪 ТЕСТ ПРЕЗЕНТАЦИИ С ПРОБЛЕМНЫМ ФАЙЛОМ")
    print("=" * 70)
    
    # Проблемный файл
    excel_path = Path(r"C:\Users\katya\Downloads\Telegram Desktop\4 сентября - четверг (2).xls")
    template_path = Path("templates/presentation_template.pptx")
    output_path = Path("test_problematic_file_presentation.pptx")
    
    if not excel_path.exists():
        print(f"❌ Excel файл не найден: {excel_path}")
        return
        
    if not template_path.exists():
        print(f"❌ Шаблон не найден: {template_path}")
        return
    
    print(f"📄 Проблемный файл: {excel_path.name}")
    print(f"📄 Шаблон: {template_path}")
    print(f"💾 Результат: {output_path}")
    print()
    
    # Сначала убедимся что извлекаем данные правильно
    print("🔍 ШАГ 1: ИЗВЛЕЧЕНИЕ РЫБНЫХ БЛЮД")
    fish_dishes = extract_fish_dishes_from_column_e(str(excel_path))
    
    if fish_dishes:
        print(f"✅ Извлечено {len(fish_dishes)} блюд:")
        for i, dish in enumerate(fish_dishes, 1):
            print(f"   {i}. '{dish.name}' | {dish.weight} | {dish.price}")
    else:
        print("❌ Рыбные блюда не извлечены")
        return
    
    # Теперь создаем презентацию
    print(f"\n🎯 ШАГ 2: СОЗДАНИЕ ПРЕЗЕНТАЦИИ")
    try:
        success, message = create_presentation_with_fish_and_side_dishes(
            str(template_path),
            str(excel_path),
            str(output_path)
        )
        
        if success:
            print(f"✅ Презентация создана: {output_path}")
            print(f"📝 Сообщение: {message}")
            
            if output_path.exists():
                size = output_path.stat().st_size
                print(f"📏 Размер файла: {size:,} байт")
                
                # Анализируем созданную презентацию
                analyze_created_presentation(str(output_path), fish_dishes)
        else:
            print(f"❌ Ошибка создания презентации: {message}")
            
    except Exception as e:
        print(f"❌ Исключение: {e}")
        import traceback
        traceback.print_exc()

def analyze_created_presentation(pptx_path: str, expected_dishes):
    """Анализирует созданную презентацию"""
    print(f"\n🔍 ШАГ 3: АНАЛИЗ СОЗДАННОЙ ПРЕЗЕНТАЦИИ")
    
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        prs = Presentation(pptx_path)
        
        if len(prs.slides) >= 6:
            slide_6 = prs.slides[5]  # 6-й слайд (индекс 5)
            print(f"🎯 Анализируем 6-й слайд:")
            
            tables_found = 0
            all_text = []
            
            for shape in slide_6.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    tables_found += 1
                    table = shape.table
                    print(f"\n📋 Таблица {tables_found}: {len(table.rows)} строк × {len(table.columns)} столбцов")
                    
                    # Показываем содержимое таблицы
                    for i, row in enumerate(table.rows):
                        row_content = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            all_text.append(cell_text)
                            if cell_text:
                                row_content.append(f"'{cell_text}'")
                            else:
                                row_content.append("'[пусто]'")
                        print(f"   Строка {i+1}: {' | '.join(row_content)}")
            
            # Проверяем, попали ли наши рыбные блюда в презентацию
            slide_text = ' '.join(all_text).upper()
            dishes_found = 0
            
            print(f"\n🎯 ПРОВЕРКА ПОПАДАНИЯ БЛЮД В ПРЕЗЕНТАЦИЮ:")
            for i, dish in enumerate(expected_dishes, 1):
                # Ищем каждое слово из названия блюда
                dish_words = dish.name.upper().split()
                words_found = sum(1 for word in dish_words if len(word) > 3 and word in slide_text)
                
                if words_found >= len(dish_words) // 2:  # Хотя бы половина слов найдена
                    dishes_found += 1
                    print(f"   ✅ Блюдо {i}: '{dish.name}' - найдено в презентации")
                else:
                    print(f"   ❌ Блюдо {i}: '{dish.name}' - НЕ найдено в презентации")
            
            print(f"\n📈 ИТОГО: {dishes_found}/{len(expected_dishes)} блюд попало в презентацию")
            
            if dishes_found == 0:
                print("❌ ПРОБЛЕМА ПОДТВЕРЖДЕНА: Названия блюд не попадают в презентацию")
                print(f"🔍 Весь текст на слайде: {slide_text[:200]}...")
            elif dishes_found == len(expected_dishes):
                print("✅ ВСЕ БЛЮДА НАЙДЕНЫ: Проблема не воспроизводится")
            else:
                print("⚠️  ЧАСТИЧНАЯ ПРОБЛЕМА: Не все блюда найдены")
                
        else:
            print(f"❌ В презентации недостаточно слайдов")
            
    except ImportError:
        print("❌ Модуль python-pptx не установлен")
    except Exception as e:
        print(f"❌ Ошибка анализа презентации: {e}")

if __name__ == "__main__":
    test_problematic_file_presentation()
