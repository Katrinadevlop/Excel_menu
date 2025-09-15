import pandas as pd
import re
import shutil
from pathlib import Path
from typing import List, Tuple, Optional
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from dataclasses import dataclass


@dataclass
class MenuItem:
    name: str
    weight: str
    price: str


def extract_dishes_from_excel_column(excel_path: str, category_keywords: List[str]) -> List[MenuItem]:
    """
    Адаптированная функция для извлечения блюд из колоночной структуры Excel.
    Работает с файлами, где категории - это заголовки колонок.
    """
    try:
        # Выбираем лист (приоритет листу с "касс")
        xls = pd.ExcelFile(excel_path)
        sheet_name = None
        for nm in xls.sheet_names:
            if 'касс' in str(nm).strip().lower():
                sheet_name = nm
                break
        if sheet_name is None and xls.sheet_names:
            sheet_name = xls.sheet_names[0]

        # Читаем весь лист
        df = pd.read_excel(excel_path, sheet_name=sheet_name, header=None, dtype=object)

        def row_text(row) -> str:
            parts = []
            for v in row:
                if pd.notna(v):
                    parts.append(str(v))
            return ' '.join(parts).strip()

        # Находим строку с заголовками колонок
        header_row = None
        category_columns = {}
        
        for i in range(min(20, len(df))):
            row_content = row_text(df.iloc[i]).upper().replace('Ё', 'Е')
            
            # Ищем строку с категориями (игнорируем регистр)
            found_categories = 0
            for keyword_set in category_keywords:
                keywords = keyword_set.upper().split(' ')
                if any(kw.upper() in row_content for kw in keywords if len(kw) > 2):
                    found_categories += 1
            
            if found_categories > 0:  # Нашли хотя бы одну категорию
                header_row = i
                
                # Определяем, в каких колонках находятся наши категории
                for col_idx, cell_value in enumerate(df.iloc[i]):
                    if pd.notna(cell_value):
                        cell_text = str(cell_value).upper().replace('Ё', 'Е')
                        
                        for keyword_set in category_keywords:
                            keywords = keyword_set.upper().split(' ')
                            if any(kw.upper() in cell_text for kw in keywords if len(kw) > 2):
                                category_columns[col_idx] = keyword_set
                                break
                break

        if header_row is None or not category_columns:
            return []

        # Извлекаем данные из колонок с нашей категорией
        dishes: List[MenuItem] = []
        
        # Начинаем с следующей строки после заголовков
        for row_idx in range(header_row + 1, len(df)):
            row = df.iloc[row_idx]
            
            # Проверяем каждую колонку с нашей категорией
            for col_idx, category in category_columns.items():
                # Проверяем, соответствует ли категория именно той, которую мы ищем
                category_matches = False
                for keyword_set in category_keywords:
                    keywords = keyword_set.upper().split(' ')
                    if any(kw.upper() in category.upper() for kw in keywords if len(kw) > 2):
                        category_matches = True
                        break
                
                if not category_matches:
                    continue
                
                # Получаем название блюда из этой колонки
                if col_idx < len(row) and pd.notna(row.iloc[col_idx]):
                    dish_name = str(row.iloc[col_idx]).strip()
                    
                    if dish_name and not dish_name.isupper() and len(dish_name) > 3:
                        # Пытаемся найти вес и цену в соседних колонках
                        weight = ""
                        price = ""
                        
                        # Проверяем следующие 2-3 колонки
                        for offset in [1, 2, 3]:
                            if col_idx + offset < len(row) and pd.notna(row.iloc[col_idx + offset]):
                                cell_value = str(row.iloc[col_idx + offset]).strip()
                                
                                # Поиск веса
                                if not weight and re.search(r'\d+.*?(?:г|шт|мл|л)', cell_value, re.IGNORECASE):
                                    weight = cell_value
                                
                                # Поиск цены
                                if not price and re.search(r'\d+', cell_value) and not re.search(r'(?:г|шт|мл|л)', cell_value):
                                    if cell_value.isdigit():
                                        price = f"{cell_value} руб."
                                    else:
                                        price = cell_value
                        
                        dishes.append(MenuItem(name=dish_name, weight=weight, price=price))
            
            # Останавливаемся, если все ячейки в строке пустые
            if not any(pd.notna(cell) for cell in row):
                break
                
        return dishes

    except Exception as e:
        print(f"Ошибка при извлечении блюд категории {category_keywords}: {e}")
        return []


def extract_dishes_from_excel(excel_path: str, category_keywords: List[str]) -> List[MenuItem]:
    """
    Универсальная функция - сначала пробуем колоночную структуру,
    если не получается - пробуем старый способ.
    """
    # Сначала пробуем колоночную структуру
    dishes = extract_dishes_from_excel_column(excel_path, category_keywords)
    if dishes:
        return dishes
    
    # Если не нашли - пробуем старый способ (строчная структура)
    return extract_dishes_from_excel_rows(excel_path, category_keywords)


def extract_dishes_from_excel_rows(excel_path: str, category_keywords: List[str]) -> List[MenuItem]:
    """
    Старая функция для строчной структуры (когда категории в отдельных строках).
    """
    try:
        # 1) Автовыбор листа (ищем по подстроке "касс", иначе первый)
        try:
            xls = pd.ExcelFile(excel_path)
            sheet_name = None
            for nm in xls.sheet_names:
                if 'касс' in str(nm).strip().lower():
                    sheet_name = nm
                    break
            if sheet_name is None and xls.sheet_names:
                sheet_name = xls.sheet_names[0]
        except Exception as e:
            sheet_name = 0

        # 2) Чтение листа без заголовков
        df = pd.read_excel(excel_path, sheet_name=sheet_name, header=None, dtype=object)

        def row_text(row) -> str:
            parts = []
            for v in row:
                if pd.notna(v):
                    parts.append(str(v))
            return ' '.join(parts).strip()

        # 3) Поиск строки начала нужной секции
        category_row = None
        for i in range(min(50, len(df))):  # Увеличим поиск до 50 строк
            s = row_text(df.iloc[i]).upper().replace('Ё', 'Е')
            if not s:
                continue
            
            # Проверяем совпадение с любым из ключевых слов категории (игнорируем регистр)
            for keyword_set in category_keywords:
                # Более гибкая проверка - частичное совпадение
                keywords = keyword_set.upper().split(' ')
                if any(kw.upper() in s for kw in keywords if len(kw) > 2):  # игнорируем короткие слова
                    category_row = i
                    break
                # Точная проверка как запасной вариант
                elif all(kw.upper() in s for kw in keywords):
                    category_row = i
                    break
            if category_row is not None:
                break

        if category_row is None:
            return []

        # 4) Хелперы распознавания категории, веса и цены
        units_pattern = r'(?:к?кал|ккал|г|гр|грамм(?:а|ов)?|мл|л|кг)'
        price_pattern = r'(?<!\\d)(\\d{1,6}(?:[\\.,]\\d{1,2})?)\\s*(?:руб\\w*|р\\.?|₽)?'

        known_cats = [
            'ЗАВТРАК', 'ПЕРВЫЕ БЛЮДА', 'ВТОРЫЕ БЛЮДА', 'ГАРНИР', 'НАПИТК', 'ДЕСЕРТ',
            'БЛЮДА ИЗ МЯСА', 'БЛЮДА ИЗ ПТИЦЫ', 'БЛЮДА ИЗ РЫБЫ', 'САЛАТЫ', 'ХОЛОДНЫЕ ЗАКУСКИ',
            'МЯСНЫЕ БЛЮДА', 'РЫБНЫЕ БЛЮДА', 'ГАРНИРЫ'
        ]

        def is_category_row(row) -> bool:
            s = row_text(row).upper()
            if not s:
                return False
            if any(k in s for k in known_cats):
                return True
            letters = ''.join(ch for ch in s if ch.isalpha())
            if letters and letters == letters.upper() and len(letters) >= 4:
                return True
            return False

        def extract_weight_from_row(row) -> str:
            for v in row:
                if pd.isna(v):
                    continue
                s = str(v)
                m = re.search(rf'(\d+[\.,]?\d*)\s*{units_pattern}', s, flags=re.IGNORECASE)
                if m:
                    qty = m.group(1).replace(',', '.')
                    unit_m = re.search(rf'{units_pattern}', s, flags=re.IGNORECASE)
                    unit = unit_m.group(0) if unit_m else ''
                    return f"{qty.replace('.', ',')} {unit}"
            return ''

        def is_weight_like(s: str) -> bool:
            return re.search(rf'{units_pattern}', s, flags=re.IGNORECASE) is not None

        def extract_price_from_row(row) -> Optional[str]:
            candidates = []
            for v in row:
                if pd.isna(v):
                    continue
                s = str(v)
                if is_weight_like(s):
                    continue
                for m in re.finditer(price_pattern, s, flags=re.IGNORECASE):
                    num = m.group(1).replace(',', '.')
                    try:
                        val = float(num)
                    except ValueError:
                        continue
                    candidates.append(val)
            if not candidates:
                return None
            val = candidates[-1]
            if abs(val - int(val)) < 1e-6:
                txt = f"{int(val)} руб."
            else:
                txt = f"{str(val).replace('.', ',')} руб."
            return txt

        # 5) Сканируем строки до следующей категории
        dishes: List[MenuItem] = []
        current_row = category_row + 1
        empty_streak = 0
        while current_row < len(df):
            row = df.iloc[current_row]
            s_join = row_text(row)

            if is_category_row(row):
                break

            if not s_join:
                empty_streak += 1
                if empty_streak >= 3:
                    break
                current_row += 1
                continue
            else:
                empty_streak = 0

            # Имя блюда — первый непустой текст в строке
            name = ''
            for v in row:
                if pd.notna(v):
                    t = str(v).strip()
                    if t:
                        name = t
                        break

            weight = extract_weight_from_row(row)
            price = extract_price_from_row(row)

            # Не требуем обязательной цены: добавляем строку по имени
            if name and not name.isupper():
                dishes.append(MenuItem(name=name, weight=weight, price=price or ""))

            current_row += 1

        return dishes

    except Exception as e:
        print(f"Ошибка при извлечении блюд категории {category_keywords}: {e}")
        return []


def extract_dishes_from_multiple_sheets(excel_path: str, sheet_names: List[str]) -> List[MenuItem]:
    """
    Извлекает все блюда из нескольких листов.
    """
    all_dishes = []
    
    try:
        xls = pd.ExcelFile(excel_path)
        
        for sheet_name in sheet_names:
            if sheet_name in xls.sheet_names:
                try:
                    df = pd.read_excel(excel_path, sheet_name=sheet_name, header=None, dtype=object)
                    
                    # Извлекаем все блюда с этого листа
                    for i in range(len(df)):
                        for j in range(len(df.columns)):
                            if pd.notna(df.iloc[i, j]):
                                dish_name = str(df.iloc[i, j]).strip()
                                
                                if (dish_name and 
                                    not dish_name.isupper() and 
                                    len(dish_name) > 3 and 
                                    not dish_name.replace(' ', '').isdigit()):
                                    
                                    # Пытаемся найти вес и цену в соседних ячейках
                                    weight = ""
                                    price = ""
                                    
                                    # Проверяем соседние ячейки
                                    for di in [-1, 0, 1]:
                                        for dj in [1, 2, 3]:
                                            try:
                                                if (i + di >= 0 and j + dj < len(df.columns) and 
                                                    i + di < len(df) and pd.notna(df.iloc[i + di, j + dj])):
                                                    cell_value = str(df.iloc[i + di, j + dj]).strip()
                                                    
                                                    # Поиск веса
                                                    if not weight and re.search(r'\d+.*?(?:г|шт|мл|л)', cell_value, re.IGNORECASE):
                                                        weight = cell_value
                                                    
                                                    # Поиск цены
                                                    if not price and re.search(r'\d+', cell_value) and not re.search(r'(?:г|шт|мл|л)', cell_value):
                                                        if cell_value.isdigit():
                                                            price = f"{cell_value} руб."
                                                        else:
                                                            price = cell_value
                                            except:
                                                continue
                                    
                                    all_dishes.append(MenuItem(name=dish_name, weight=weight, price=price))
                                    
                except Exception as e:
                    print(f"Ошибка при чтении листа {sheet_name}: {e}")
                    
    except Exception as e:
        print(f"Ошибка при открытии файла: {e}")
    
    return all_dishes


def _upper_no_yo(s: str) -> str:
    return s.upper().replace('Ё', 'Е') if isinstance(s, str) else str(s).upper().replace('Ё', 'Е')


def extract_dishes_from_excel_rows_with_stop(excel_path: str, category_keywords: List[str], stop_keywords: List[str]) -> List[MenuItem]:
    """
    Вариант построчного извлечения, который останавливается, когда встречается
    новая категория из stop_keywords (например, "СЭНДВИЧИ").
    """
    try:
        xls = pd.ExcelFile(excel_path)
        sheet_name = None
        for nm in xls.sheet_names:
            if 'касс' in str(nm).strip().lower():
                sheet_name = nm
                break
        if sheet_name is None and xls.sheet_names:
            sheet_name = xls.sheet_names[0]

        df = pd.read_excel(excel_path, sheet_name=sheet_name, header=None, dtype=object)

        def row_text(row) -> str:
            parts = []
            for v in row:
                if pd.notna(v):
                    parts.append(str(v))
            return ' '.join(parts).strip()

        # 1) Находим стартовую строку по ключевым словам категории
        header_row = None
        for i in range(min(80, len(df))):
            s = _upper_no_yo(row_text(df.iloc[i]))
            if not s:
                continue
            for keyword_set in category_keywords:
                keywords = _upper_no_yo(keyword_set).split(' ')
                if any(kw and kw in s for kw in keywords if len(kw) > 2):
                    header_row = i
                    break
            if header_row is not None:
                break

        if header_row is None:
            return []

        # 2) Подготовим служебные функции под вес и цену (повторно используем из старой функции)
        units_pattern = r'(?:к?кал|ккал|г|гр|грамм(?:а|ов)?|мл|л|кг)'
        price_pattern = r'(?<!\\d)(\\d{1,6}(?:[\\.,]\\d{1,2})?)\\s*(?:руб\\w*|р\\.?|₽)?'

        def is_category_row(row) -> bool:
            s = _upper_no_yo(row_text(row))
            if not s:
                return False
            letters = ''.join(ch for ch in s if ch.isalpha())
            if letters and letters == letters.upper() and len(letters) >= 4:
                return True
            return False

        def extract_weight_from_row(row) -> str:
            for v in row:
                if pd.isna(v):
                    continue
                s = str(v)
                m = re.search(rf'(\\d+[\\.,]?\\d*)\\s*{units_pattern}', s, flags=re.IGNORECASE)
                if m:
                    qty = m.group(1).replace(',', '.')
                    unit_m = re.search(rf'{units_pattern}', s, flags=re.IGNORECASE)
                    unit = unit_m.group(0) if unit_m else ''
                    return f"{qty.replace('.', ',')} {unit}"
            return ''

        def is_weight_like(s: str) -> bool:
            return re.search(rf'{units_pattern}', s, flags=re.IGNORECASE) is not None

        def extract_price_from_row(row) -> Optional[str]:
            candidates = []
            for v in row:
                if pd.isna(v):
                    continue
                s = str(v)
                if is_weight_like(s):
                    continue
                for m in re.finditer(price_pattern, s, flags=re.IGNORECASE):
                    num = m.group(1).replace(',', '.')
                    try:
                        val = float(num)
                    except ValueError:
                        continue
                    candidates.append(val)
            if not candidates:
                return None
            val = candidates[-1]
            if abs(val - int(val)) < 1e-6:
                txt = f"{int(val)} руб."
            else:
                txt = f"{str(val).replace('.', ',')} руб."
            return txt

        # 3) Сканируем строки до ближайшей категории из stop_keywords или любой новой категории
        dishes: List[MenuItem] = []
        current_row = header_row + 1
        empty_streak = 0
        stop_upper = [_upper_no_yo(x) for x in stop_keywords]

        while current_row < len(df):
            row = df.iloc[current_row]
            s_join = row_text(row)
            s_upper = _upper_no_yo(s_join)

            # Стоп по встрече новой категории «СЭНДВИЧИ/СЕНДВИЧИ»
            if is_category_row(row) and any(st in s_upper for st in stop_upper):
                break

            # Остановка при любой новой категории (кроме строк обычных блюд)
            if is_category_row(row) and s_join:
                # Если это не "пустая" строка блюда — прекращаем сбор текущей секции
                break

            if not s_join:
                empty_streak += 1
                if empty_streak >= 3:
                    break
                current_row += 1
                continue
            else:
                empty_streak = 0

            # Имя блюда — первый текст в строке
            name = ''
            for v in row:
                if pd.notna(v):
                    t = str(v).strip()
                    if t:
                        name = t
                        break

            weight = extract_weight_from_row(row)
            price = extract_price_from_row(row)

            if name and not name.isupper():
                dishes.append(MenuItem(name=name, weight=weight, price=price or ""))

            current_row += 1

        return dishes

    except Exception as e:
        print(f"Ошибка при выборочном извлечении (до стоп-категории) {category_keywords}: {e}")
        return []


def extract_salads_from_excel(excel_path: str) -> List[MenuItem]:
    """Извлекает салаты и холодные закуски от заголовка 'САЛАТЫ И ХОЛОДНЫЕ ЗАКУСКИ' до 'СЭНДВИЧИ'."""
    try:
        # Сначала пробуем найти по точному диапазону от заголовка до СЭНДВИЧИ
        salads = extract_salads_by_range(excel_path)
        if salads:
            return salads
        
        # Если не получилось - пробуем старый способ через листы
        return extract_dishes_from_multiple_sheets(excel_path, ['Хц', 'Холодные', 'Салаты', 'касса '])
    except Exception as e:
        print(f"Ошибка при извлечении салатов: {e}")
        return []


def extract_salads_by_range(excel_path: str) -> List[MenuItem]:
    """Извлекает салаты из диапазона от 'САЛАТЫ И ХОЛОДНЫЕ ЗАКУСКИ' до 'СЭНДВИЧИ'."""
    try:
        # Выбираем лист (приоритет листу с "касс")
        xls = pd.ExcelFile(excel_path)
        sheet_name = None
        for nm in xls.sheet_names:
            if 'касс' in str(nm).strip().lower():
                sheet_name = nm
                break
        if sheet_name is None and xls.sheet_names:
            sheet_name = xls.sheet_names[0]

        # Читаем весь лист
        df = pd.read_excel(excel_path, sheet_name=sheet_name, header=None, dtype=object)
        
        def row_text(row) -> str:
            parts = []
            for v in row:
                if pd.notna(v):
                    parts.append(str(v))
            return ' '.join(parts).strip()
        
        # Находим строку с заголовком "САЛАТЫ И ХОЛОДНЫЕ ЗАКУСКИ"
        start_row = None
        end_row = None
        
        for i in range(len(df)):
            row_content = row_text(df.iloc[i]).upper().replace('Ё', 'Е')
            
            # Ищем начало секции салатов
            if start_row is None:
                if ('САЛАТ' in row_content and 'ХОЛОДН' in row_content and 'ЗАКУСК' in row_content) or \
                   ('САЛАТЫ' in row_content and ('ХОЛОДНЫЕ' in row_content or 'ЗАКУСКИ' in row_content)):
                    start_row = i
                    print(f"Найден заголовок салатов в строке {i + 1}: {row_content}")
                    continue
            
            # Ищем конец секции (СЭНДВИЧИ)
            if start_row is not None and end_row is None:
                if 'СЭНДВИЧ' in row_content:
                    end_row = i
                    print(f"Найден заголовок сэндвичей в строке {i + 1}: {row_content}")
                    break
        
        if start_row is None:
            print("Заголовок 'САЛАТЫ И ХОЛОДНЫЕ ЗАКУСКИ' не найден")
            return []
        
        if end_row is None:
            # Если не найдены СЭНДВИЧИ, берем до конца файла или до следующей крупной категории
            end_row = len(df)
            for i in range(start_row + 1, len(df)):
                row_content = row_text(df.iloc[i]).upper().replace('Ё', 'Е')
                # Ищем другие крупные категории как стоп-слова
                if any(category in row_content for category in [
                    'ПЕРВЫЕ БЛЮДА', 'ВТОРЫЕ БЛЮДА', 'ГОРЯЧИЕ БЛЮДА', 
                    'МЯСНЫЕ БЛЮДА', 'РЫБНЫЕ БЛЮДА', 'ГАРНИРЫ', 'НАПИТКИ'
                ]):
                    end_row = i
                    break
        
        print(f"Обрабатываем салаты от строки {start_row + 1} до строки {end_row}")
        
        # Извлекаем блюда из найденного диапазона
        dishes: List[MenuItem] = []
        
        for i in range(start_row + 1, end_row):
            if i >= len(df):
                break
                
            row = df.iloc[i]
            row_content = row_text(row)
            
            # Пропускаем пустые строки
            if not row_content.strip():
                continue
                
            # Пропускаем строки, которые выглядят как подзаголовки (все заглавные)
            if row_content.isupper() and len(row_content) > 10:
                continue
                
            # Ищем название блюда в первой непустой ячейке
            dish_name = ""
            dish_weight = ""
            dish_price = ""
            
            # Проходим по всем ячейкам строки
            for j, cell_value in enumerate(row):
                if pd.notna(cell_value):
                    cell_text = str(cell_value).strip()
                    if not cell_text:
                        continue
                    
                    # Первая значимая ячейка - это название блюда
                    if not dish_name and not cell_text.isupper() and len(cell_text) > 3:
                        # Проверяем, что это не цена и не вес
                        if not re.match(r'^\d+([.,]\d+)?\s*(руб|₽|р\.?)?$', cell_text) and \
                           not re.search(r'\d+\s*(г|гр|мл|л|шт)', cell_text, re.IGNORECASE):
                            dish_name = cell_text
                            continue
                    
                    # Ищем вес (содержит единицы измерения)
                    if not dish_weight and re.search(r'\d+.*?(г|гр|грамм|мл|л|кг|шт)', cell_text, re.IGNORECASE):
                        dish_weight = cell_text
                        continue
                    
                    # Ищем цену (число, возможно с "руб" или символом рубля)
                    if not dish_price and re.search(r'\d+', cell_text):
                        # Проверяем, что это не вес
                        if not re.search(r'г|гр|грамм|мл|л|кг|шт', cell_text, re.IGNORECASE):
                            if cell_text.replace('.', '').replace(',', '').isdigit():
                                dish_price = f"{cell_text} руб."
                            elif re.search(r'\d+.*?(руб|₽|р\.?)', cell_text, re.IGNORECASE):
                                dish_price = cell_text
                            else:
                                # Просто число - добавляем "руб."
                                number_match = re.search(r'\d+([.,]\d+)?', cell_text)
                                if number_match:
                                    dish_price = f"{number_match.group()} руб."
            
            # Если нашли название блюда, добавляем его
            if dish_name and len(dish_name) > 2:
                dishes.append(MenuItem(name=dish_name, weight=dish_weight, price=dish_price))
                print(f"Найден салат: {dish_name} | {dish_weight} | {dish_price}")
        
        print(f"Всего найдено салатов: {len(dishes)}")
        return dishes
        
    except Exception as e:
        print(f"Ошибка при извлечении салатов по диапазону: {e}")
        return []


def extract_first_courses_from_excel(excel_path: str) -> List[MenuItem]:
    """Извлекает первые блюда - только супы."""
    # Пробуем найти в основном листе (касса) - у вас там ПЕРВЫЕ БЛЮДА
    keywords = ['ПЕРВЫЕ БЛЮДА', 'ПЕРВЫЕ']
    all_dishes = extract_dishes_from_excel(excel_path, keywords)
    
    # Фильтруем - оставляем только супы
    soups = []
    for dish in all_dishes:
        dish_name_lower = dish.name.lower()
        # Проверяем, что это суп или бульон
        if any(soup_word in dish_name_lower for soup_word in [
            'суп', 'бульон', 'солянка', 'борщ', 'щи', 'окрошка', 'харчо', 'рассольник',
            'крем-суп', 'суп-пюре', 'овощной суп', 'рыбный суп', 'куриный суп'
        ]):
            soups.append(dish)
    
    print(f"Найдено {len(all_dishes)} блюд в секции 'Первые блюда', из них {len(soups)} супов")
    return soups


def extract_meat_dishes_by_range(excel_path: str) -> List[MenuItem]:
    """Извлекает блюда из мяса из точного диапазона от 'БЛЮДА ИЗ МЯСА' до 'БЛЮДА ИЗ ПТИЦЫ'."""
    try:
        # Выбираем лист (приоритет листу с "касс")
        xls = pd.ExcelFile(excel_path)
        sheet_name = None
        for nm in xls.sheet_names:
            if 'касс' in str(nm).strip().lower():
                sheet_name = nm
                break
        if sheet_name is None and xls.sheet_names:
            sheet_name = xls.sheet_names[0]

        # Читаем весь лист
        df = pd.read_excel(excel_path, sheet_name=sheet_name, header=None, dtype=object)
        
        def row_text(row) -> str:
            parts = []
            for v in row:
                if pd.notna(v):
                    parts.append(str(v))
            return ' '.join(parts).strip()
        
        # Находим строку с заголовком "БЛЮДА ИЗ МЯСА"
        start_row = None
        end_row = None
        
        for i in range(len(df)):
            row_content = row_text(df.iloc[i]).upper().replace('Ё', 'Е')
            
            # Ищем начало секции мясных блюд
            if start_row is None:
                if 'БЛЮДА ИЗ МЯСА' in row_content or 'МЯСНЫЕ БЛЮДА' in row_content:
                    start_row = i
                    print(f"Найден заголовок мясных блюд в строке {i + 1}: {row_content}")
                    continue
            
            # Ищем конец секции (БЛЮДА ИЗ ПТИЦЫ)
            if start_row is not None and end_row is None:
                if 'БЛЮДА ИЗ ПТИЦЫ' in row_content or ('ПТИЦ' in row_content and 'БЛЮДА' in row_content):
                    end_row = i
                    print(f"Найден заголовок блюд из птицы в строке {i + 1}: {row_content}")
                    break
        
        if start_row is None:
            print("Заголовок 'БЛЮДА ИЗ МЯСА' не найден")
            return []
        
        if end_row is None:
            # Если не найдены блюда из птицы, ищем другие стоп-категории
            end_row = len(df)
            for i in range(start_row + 1, len(df)):
                row_content = row_text(df.iloc[i]).upper().replace('Ё', 'Е')
                # Ищем другие категории как стоп-слова
                if any(category in row_content for category in [
                    'РЫБНЫЕ БЛЮДА', 'БЛЮДА ИЗ РЫБЫ', 'ГАРНИРЫ', 'НАПИТКИ', 
                    'ДЕСЕРТЫ', 'САЛАТЫ', 'СЭНДВИЧ'
                ]):
                    end_row = i
                    break
        
        print(f"Обрабатываем мясные блюда от строки {start_row + 1} до строки {end_row}")
        
        # Извлекаем блюда из найденного диапазона
        dishes: List[MenuItem] = []
        
        for i in range(start_row + 1, end_row):
            if i >= len(df):
                break
                
            row = df.iloc[i]
            row_content = row_text(row)
            
            # Пропускаем пустые строки
            if not row_content.strip():
                continue
                
            # Пропускаем строки, которые выглядят как подзаголовки (все заглавные)
            if row_content.isupper() and len(row_content) > 10:
                continue
                
            # Ищем название блюда в первой непустой ячейке
            dish_name = ""
            dish_weight = ""
            dish_price = ""
            
            # Проходим по всем ячейкам строки
            for j, cell_value in enumerate(row):
                if pd.notna(cell_value):
                    cell_text = str(cell_value).strip()
                    if not cell_text:
                        continue
                    
                    # Первая значимая ячейка - это название блюда
                    if not dish_name and not cell_text.isupper() and len(cell_text) > 3:
                        # Проверяем, что это не цена и не вес
                        if not re.match(r'^\d+([.,]\d+)?\s*(руб|₽|р\.?)?$', cell_text) and \
                           not re.search(r'\d+\s*(г|гр|мл|л|шт)', cell_text, re.IGNORECASE):
                            dish_name = cell_text
                            continue
                    
                    # Ищем вес (содержит единицы измерения)
                    if not dish_weight and re.search(r'\d+.*?(г|гр|грамм|мл|л|кг|шт)', cell_text, re.IGNORECASE):
                        dish_weight = cell_text
                        continue
                    
                    # Ищем цену (число, возможно с "руб" или символом рубля)
                    if not dish_price and re.search(r'\d+', cell_text):
                        # Проверяем, что это не вес
                        if not re.search(r'г|гр|грамм|мл|л|кг|шт', cell_text, re.IGNORECASE):
                            if cell_text.replace('.', '').replace(',', '').isdigit():
                                dish_price = f"{cell_text} руб."
                            elif re.search(r'\d+.*?(руб|₽|р\.?)', cell_text, re.IGNORECASE):
                                dish_price = cell_text
                            else:
                                # Просто число - добавляем "руб."
                                number_match = re.search(r'\d+([.,]\d+)?', cell_text)
                                if number_match:
                                    dish_price = f"{number_match.group()} руб."
            
            # Если нашли название блюда, добавляем его в список
            # Не нужна дополнительная фильтрация, так как мы уже в разделе "БЛЮДА ИЗ МЯСА"
            if dish_name and len(dish_name) > 2:
                dishes.append(MenuItem(name=dish_name, weight=dish_weight, price=dish_price))
                print(f"Найдено блюдо из мяса: {dish_name} | {dish_weight} | {dish_price}")
        
        print(f"Всего найдено мясных блюд: {len(dishes)}")
        return dishes
        
    except Exception as e:
        print(f"Ошибка при извлечении мясных блюд по диапазону: {e}")
        return []


def extract_meat_dishes_from_excel(excel_path: str) -> List[MenuItem]:
    """Извлекает блюда из мяса - точный диапазон от 'БЛЮДА ИЗ МЯСА' до 'БЛЮДА ИЗ ПТИЦЫ'."""
    try:
        # Пробуем найти точный диапазон от "БЛЮДА ИЗ МЯСА" до "БЛЮДА ИЗ ПТИЦЫ"
        dishes = extract_meat_dishes_by_range(excel_path)
        if dishes:
            return dishes
            
        # Если не нашли по диапазону, пробуем с остановкой
        stop_keywords = ['БЛЮДА ИЗ ПТИЦЫ', 'ПТИЦА', 'РЫБНЫЕ БЛЮДА', 'БЛЮДА ИЗ РЫБЫ']
        keywords = ['БЛЮДА ИЗ МЯСА', 'МЯСНЫЕ БЛЮДА', 'МЯСО']
        
        dishes = extract_dishes_from_excel_rows_with_stop(excel_path, keywords, stop_keywords)
        
        if dishes:
            print(f"Найдено {len(dishes)} мясных блюд через построчный поиск с остановкой")
            return dishes
            
        # Последняя попытка - поиск по ключевым словам с фильтрацией
        print("Поиск мясных блюд через колоночную структуру...")
        all_keywords = ['БЛЮДА ИЗ МЯСА', 'МЯСНЫЕ БЛЮДА', 'МЯСО']
        dishes = extract_dishes_from_excel(excel_path, all_keywords)
        
        # Фильтруем, исключая супы и первые блюда
        if dishes:
            filtered_dishes = []
            for dish in dishes:
                dish_name_lower = dish.name.lower()
                # Исключаем супы, бульоны и явно первые блюда
                if not any(soup_word in dish_name_lower for soup_word in [
                    'суп', 'бульон', 'солянка', 'борщ', 'щи', 'окрошка', 'харчо'
                ]):
                    filtered_dishes.append(dish)
            
            if filtered_dishes:
                return filtered_dishes
        
        return []
        
    except Exception as e:
        print(f"Ошибка при извлечении мясных блюд: {e}")
        return []


def extract_poultry_dishes_by_range(excel_path: str) -> List[MenuItem]:
    """Извлекает блюда из птицы из точного диапазона от 'БЛЮДА ИЗ ПТИЦЫ' до 'БЛЮДА ИЗ РЫБЫ'."""
    try:
        # Выбираем лист (приоритет листу с "касс")
        xls = pd.ExcelFile(excel_path)
        sheet_name = None
        for nm in xls.sheet_names:
            if 'касс' in str(nm).strip().lower():
                sheet_name = nm
                break
        if sheet_name is None and xls.sheet_names:
            sheet_name = xls.sheet_names[0]

        # Читаем весь лист
        df = pd.read_excel(excel_path, sheet_name=sheet_name, header=None, dtype=object)
        
        def row_text(row) -> str:
            parts = []
            for v in row:
                if pd.notna(v):
                    parts.append(str(v))
            return ' '.join(parts).strip()
        
        # Находим строку с заголовком "БЛЮДА ИЗ ПТИЦЫ"
        start_row = None
        end_row = None
        
        for i in range(len(df)):
            row_content = row_text(df.iloc[i]).upper().replace('Ё', 'Е')
            
            # Ищем начало секции блюд из птицы
            if start_row is None:
                if 'БЛЮДА ИЗ ПТИЦЫ' in row_content or ('ПТИЦ' in row_content and 'БЛЮДА' in row_content):
                    start_row = i
                    print(f"Найден заголовок блюд из птицы в строке {i + 1}: {row_content}")
                    continue
            
            # Ищем конец секции (БЛЮДА ИЗ РЫБЫ)
            if start_row is not None and end_row is None:
                if 'БЛЮДА ИЗ РЫБЫ' in row_content or ('РЫБ' in row_content and 'БЛЮДА' in row_content):
                    end_row = i
                    print(f"Найден заголовок блюд из рыбы в строке {i + 1}: {row_content}")
                    break
        
        if start_row is None:
            print("Заголовок 'БЛЮДА ИЗ ПТИЦЫ' не найден")
            return []
        
        if end_row is None:
            # Если не найдены блюда из рыбы, ищем другие стоп-категории
            end_row = len(df)
            for i in range(start_row + 1, len(df)):
                row_content = row_text(df.iloc[i]).upper().replace('Ё', 'Е')
                # Ищем другие категории как стоп-слова
                if any(category in row_content for category in [
                    'ГАРНИРЫ', 'НАПИТКИ', 'ДЕСЕРТЫ', 'САЛАТЫ', 'СЭНДВИЧ'
                ]):
                    end_row = i
                    break
        
        print(f"Обрабатываем блюда из птицы от строки {start_row + 1} до строки {end_row}")
        
        # Извлекаем блюда из найденного диапазона
        dishes: List[MenuItem] = []
        
        for i in range(start_row + 1, end_row):
            if i >= len(df):
                break
                
            row = df.iloc[i]
            row_content = row_text(row)
            
            # Пропускаем пустые строки
            if not row_content.strip():
                continue
                
            # Пропускаем строки, которые выглядят как подзаголовки (все заглавные)
            if row_content.isupper() and len(row_content) > 10:
                continue
                
            # Ищем название блюда в первой непустой ячейке
            dish_name = ""
            dish_weight = ""
            dish_price = ""
            
            # Проходим по всем ячейкам строки
            for j, cell_value in enumerate(row):
                if pd.notna(cell_value):
                    cell_text = str(cell_value).strip()
                    if not cell_text:
                        continue
                    
                    # Первая значимая ячейка - это название блюда
                    if not dish_name and not cell_text.isupper() and len(cell_text) > 3:
                        # Проверяем, что это не цена и не вес
                        if not re.match(r'^\d+([.,]\d+)?\s*(руб|₽|р\.?)?$', cell_text) and \
                           not re.search(r'\d+\s*(г|гр|мл|л|шт)', cell_text, re.IGNORECASE):
                            dish_name = cell_text
                            continue
                    
                    # Ищем вес (содержит единицы измерения)
                    if not dish_weight and re.search(r'\d+.*?(г|гр|грамм|мл|л|кг|шт)', cell_text, re.IGNORECASE):
                        dish_weight = cell_text
                        continue
                    
                    # Ищем цену (число, возможно с "руб" или символом рубля)
                    if not dish_price and re.search(r'\d+', cell_text):
                        # Проверяем, что это не вес
                        if not re.search(r'г|гр|грамм|мл|л|кг|шт', cell_text, re.IGNORECASE):
                            if cell_text.replace('.', '').replace(',', '').isdigit():
                                dish_price = f"{cell_text} руб."
                            elif re.search(r'\d+.*?(руб|₽|р\.?)', cell_text, re.IGNORECASE):
                                dish_price = cell_text
                            else:
                                # Просто число - добавляем "руб."
                                number_match = re.search(r'\d+([.,]\d+)?', cell_text)
                                if number_match:
                                    dish_price = f"{number_match.group()} руб."
            
            # Если нашли название блюда, добавляем его в список
            if dish_name and len(dish_name) > 2:
                dishes.append(MenuItem(name=dish_name, weight=dish_weight, price=dish_price))
                print(f"Найдено блюдо из птицы: {dish_name} | {dish_weight} | {dish_price}")
        
        print(f"Всего найдено блюд из птицы: {len(dishes)}")
        return dishes
        
    except Exception as e:
        print(f"Ошибка при извлечении блюд из птицы по диапазону: {e}")
        return []


def extract_poultry_dishes_from_excel(excel_path: str) -> List[MenuItem]:
    """Извлекает блюда из птицы - точный диапазон от 'БЛЮДА ИЗ ПТИЦЫ' до 'БЛЮДА ИЗ РЫБЫ'."""
    try:
        # Пробуем найти точный диапазон от "БЛЮДА ИЗ ПТИЦЫ" до "БЛЮДА ИЗ РЫБЫ"
        dishes = extract_poultry_dishes_by_range(excel_path)
        if dishes:
            return dishes
            
        # Если не нашли по диапазону, пробуем с остановкой
        stop_keywords = ['БЛЮДА ИЗ РЫБЫ', 'РЫБА', 'РЫБНЫЕ БЛЮДА', 'ГАРНИРЫ']
        keywords = ['БЛЮДА ИЗ ПТИЦЫ', 'ПТИЦА', 'КУРИНЫЕ БЛЮДА']
        
        dishes = extract_dishes_from_excel_rows_with_stop(excel_path, keywords, stop_keywords)
        
        if dishes:
            print(f"Найдено {len(dishes)} блюд из птицы через построчный поиск с остановкой")
            return dishes
            
        # Последняя попытка - поиск в листах
        print("Поиск блюд из птицы через листы...")
        return extract_dishes_from_multiple_sheets(excel_path, ['Обед', 'Гц', 'Птица', 'касса '])
        
    except Exception as e:
        print(f"Ошибка при извлечении блюд из птицы: {e}")
        return []


def extract_fish_dishes_from_excel(excel_path: str) -> List[MenuItem]:
    """Извлекает блюда из рыбы."""
    # Пробуем найти в листах Обед, Гц
    return extract_dishes_from_multiple_sheets(excel_path, ['Обед', 'Гц', 'Рыба', 'касса '])


def extract_side_dishes_from_excel(excel_path: str) -> List[MenuItem]:
    """Извлекает гарниры."""
    # Пробуем найти в листах Раздача, Обед, Гц
    return extract_dishes_from_multiple_sheets(excel_path, ['Раздача', 'Обед', 'Гц', 'Гарниры', 'касса '])


def update_slide_with_dishes(slide, dishes: List[MenuItem]) -> bool:
    """
    Обновляет один слайд презентации, вставляя данные блюд в таблицу
    с форматированием Gilroy Medium 28pt, белый цвет, автоуменьшение и отступом 10 пикселей.
    """
    try:
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        # Найдем таблицу на слайде
        table_shape = None
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_shape = shape
                break
                
        if table_shape is None:
            return False
            
        table = table_shape.table
        
        # Получаем количество строк в таблице
        total_rows = len(table.rows)
        
        # Определяем оптимальный размер шрифта (начинаем с 28pt)
        available_rows = total_rows - 1  # -1 для заголовка
        dishes_to_show = len(dishes)
        
        # Начинаем с 28pt и уменьшаем, если нужно
        if dishes_to_show <= available_rows:
            font_size = 28  # Оптимальный размер
        elif dishes_to_show <= available_rows * 1.5:
            font_size = 24  # Немного уменьшаем
        elif dishes_to_show <= available_rows * 2:
            font_size = 20  # Еще уменьшаем
        elif dishes_to_show <= available_rows * 3:
            font_size = 16  # Минимально читаемый
        else:
            font_size = 14  # Критически маленький
        
        # Ограничиваем количество блюд доступным местом
        dishes_to_fill = dishes[:available_rows]
        
        # Очищаем все строки кроме первой (заголовки) и заполняем их блюдами
        for i, dish in enumerate(dishes_to_fill):
            row_idx = i + 1  # +1 потому что 0 - это заголовок
            
            if row_idx < total_rows:
                row = table.rows[row_idx]
                
                # Заполняем и форматируем ячейки
                if len(row.cells) >= 3:
                    # Название блюда
                    cell_name = row.cells[0]
                    cell_name.text = dish.name
                    if cell_name.text_frame.paragraphs:
                        paragraph = cell_name.text_frame.paragraphs[0]
                        paragraph.alignment = PP_ALIGN.LEFT
                        # Устанавливаем отступ в 10 пикселей
                        cell_name.text_frame.margin_left = Pt(10)
                        cell_name.text_frame.margin_right = Pt(10)
                        cell_name.text_frame.margin_top = Pt(10)
                        cell_name.text_frame.margin_bottom = Pt(10)
                        if paragraph.runs:
                            run = paragraph.runs[0]
                            run.font.name = 'Gilroy Medium'
                            run.font.size = Pt(font_size)
                            run.font.color.rgb = RGBColor(255, 255, 255)  # Белый цвет
                    
                    # Вес/объем
                    cell_weight = row.cells[1]
                    cell_weight.text = dish.weight
                    if cell_weight.text_frame.paragraphs:
                        paragraph = cell_weight.text_frame.paragraphs[0]
                        paragraph.alignment = PP_ALIGN.CENTER
                        # Устанавливаем отступ в 10 пикселей
                        cell_weight.text_frame.margin_left = Pt(10)
                        cell_weight.text_frame.margin_right = Pt(10)
                        cell_weight.text_frame.margin_top = Pt(10)
                        cell_weight.text_frame.margin_bottom = Pt(10)
                        if paragraph.runs:
                            run = paragraph.runs[0]
                            run.font.name = 'Gilroy Medium'
                            run.font.size = Pt(font_size)
                            run.font.color.rgb = RGBColor(255, 255, 255)
                    
                    # Цена
                    cell_price = row.cells[2]
                    cell_price.text = dish.price
                    if cell_price.text_frame.paragraphs:
                        paragraph = cell_price.text_frame.paragraphs[0]
                        paragraph.alignment = PP_ALIGN.CENTER
                        # Устанавливаем отступ в 10 пикселей
                        cell_price.text_frame.margin_left = Pt(10)
                        cell_price.text_frame.margin_right = Pt(10)
                        cell_price.text_frame.margin_top = Pt(10)
                        cell_price.text_frame.margin_bottom = Pt(10)
                        if paragraph.runs:
                            run = paragraph.runs[0]
                            run.font.name = 'Gilroy Medium'
                            run.font.size = Pt(font_size)
                            run.font.color.rgb = RGBColor(255, 255, 255)
        
        # Очищаем оставшиеся строки
        for i in range(len(dishes_to_fill) + 1, total_rows):
            if i < len(table.rows):
                row = table.rows[i]
                for j in range(len(row.cells)):
                    row.cells[j].text = ""
                
        return True
        
    except Exception as e:
        print(f"Ошибка при обновлении слайда: {e}")
        return False


def update_presentation_with_all_categories(presentation_path: str, all_dishes: dict, output_path: str) -> bool:
    """
    Обновляет презентацию, вставляя данные всех категорий блюд в соответствующие слайды.
    
    Args:
        presentation_path: путь к шаблону презентации
        all_dishes: словарь с данными блюд для каждой категории
                   {'salads': [...], 'first_courses': [...], 'meat': [...], 
                    'poultry': [...], 'fish': [...], 'side_dishes': [...]}
        output_path: путь для сохранения готовой презентации
    """
    try:
        # Копируем исходную презентацию
        shutil.copy2(presentation_path, output_path)
        
        # Открываем презентацию
        prs = Presentation(output_path)
        
        # Проверяем, что у нас достаточно слайдов
        if len(prs.slides) < 7:
            return False
            
        # Маппинг категорий на индексы слайдов
        slide_mapping = {
            'salads': 1,        # Слайд 2: САЛАТЫ И ХОЛОДНЫЕ ЗАКУСКИ
            'first_courses': 2, # Слайд 3: ПЕРВЫЕ БЛЮДА
            'meat': 3,          # Слайд 4: МЯСНЫЕ БЛЮДА
            'poultry': 4,       # Слайд 5: БЛЮДА ИЗ ПТИЦЫ
            'fish': 5,          # Слайд 6: РЫБНЫЕ БЛЮДА
            'side_dishes': 6    # Слайд 7: ГАРНИРЫ
        }
        
        # Обновляем каждый слайд соответствующими данными
        success_count = 0
        for category, slide_idx in slide_mapping.items():
            if category in all_dishes and all_dishes[category]:
                slide = prs.slides[slide_idx]
                if update_slide_with_dishes(slide, all_dishes[category]):
                    success_count += 1
                    print(f"Слайд {slide_idx + 1} ({category}): добавлено {len(all_dishes[category])} блюд")
                else:
                    print(f"Ошибка при обновлении слайда {slide_idx + 1} ({category})")
            else:
                print(f"Нет данных для категории {category}")
                
        # Сохраняем презентацию
        prs.save(output_path)
        
        return success_count > 0
        
    except Exception as e:
        print(f"Ошибка при обновлении презентации: {e}")
        return False


def update_presentation_with_salads(presentation_path: str, salads: List[MenuItem], output_path: str) -> bool:
    """
    Обновляет презентацию, вставляя данные салатов во второй слайд.
    (Оставлена для совместимости со старым кодом)
    """
    all_dishes = {'salads': salads}
    return update_presentation_with_all_categories(presentation_path, all_dishes, output_path)


def create_presentation_with_excel_data(template_path: str, excel_path: str, output_path: str) -> Tuple[bool, str]:
    """
    Создает презентацию с салатами, первыми блюдами, блюдами из мяса и птицы.
    Остальные слайды остаются пустыми.

    Returns:
        Tuple[bool, str]: (успех, сообщение)
    """
    try:
        # Проверяем существование файлов
        if not Path(template_path).exists():
            return False, f"Шаблон презентации не найден: {template_path}"
            
        if not Path(excel_path).exists():
            return False, f"Excel файл не найден: {excel_path}"
        
        # Извлекаем салаты из Excel
        print(f"🔍 Ищем салаты в файле: {excel_path}")
        salads = extract_salads_from_excel(excel_path)
        print(f"Салаты: найдено {len(salads)} блюд")
        
        # Если салаты не найдены, пробуем альтернативные способы поиска
        if len(salads) == 0:
            keywords = ['САЛАТЫ', 'ХОЛОДНЫЕ ЗАКУСКИ', 'САЛАТЫ И ХОЛОДНЫЕ ЗАКУСКИ']
            salads = extract_dishes_from_excel(excel_path, keywords)
            print(f"Салаты (альтернативный поиск): найдено {len(salads)} блюд")
        
        # Извлекаем первые блюда
        print(f"🔍 Ищем первые блюда в файле: {excel_path}")
        first_courses = extract_first_courses_from_excel(excel_path)
        print(f"Первые блюда: найдено {len(first_courses)} блюд")
        
        # Извлекаем блюда из мяса
        print(f"🔍 Ищем блюда из мяса в файле: {excel_path}")
        meat_dishes = extract_meat_dishes_from_excel(excel_path)
        print(f"Блюда из мяса: найдено {len(meat_dishes)} блюд")
        
        # Извлекаем блюда из птицы
        print(f"🔍 Ищем блюда из птицы в файле: {excel_path}")
        poultry_dishes = extract_poultry_dishes_from_excel(excel_path)
        print(f"Блюда из птицы: найдено {len(poultry_dishes)} блюд")
        
        # Проверяем, что хотя бы одна категория найдена
        total_dishes = len(salads) + len(first_courses) + len(meat_dishes) + len(poultry_dishes)
        
        if total_dishes == 0:
            # Попробуем показать содержимое файла для диагностики
            try:
                import pandas as pd
                xls = pd.ExcelFile(excel_path)
                print(f"📋 Листы в файле: {xls.sheet_names}")
                
                sheet_name = xls.sheet_names[0]
                if 'касс' in str(sheet_name).lower():
                    sheet_name = next((nm for nm in xls.sheet_names if 'касс' in str(nm).lower()), xls.sheet_names[0])
                
                df = pd.read_excel(excel_path, sheet_name=sheet_name, header=None, dtype=object)
                print(f"📊 Размер данных: {len(df)} строк")
                print("📝 Первые 20 строк с содержимым:")
                
                def row_text(row) -> str:
                    parts = []
                    for v in row:
                        if pd.notna(v):
                            parts.append(str(v))
                    return ' '.join(parts).strip()
                
                for i in range(min(20, len(df))):
                    content = row_text(df.iloc[i])
                    if content.strip():
                        print(f"  Строка {i+1}: {content[:100]}")
                        
            except Exception as diag_error:
                print(f"Ошибка диагностики: {diag_error}")
                
            return False, f"В Excel файле не найдены блюда указанных категорий. Проверьте структуру файла и названия категорий."

        # Создаем словарь с найденными блюдами
        all_dishes = {
            'salads': salads,
            'first_courses': first_courses,
            'meat': meat_dishes,
            'poultry': poultry_dishes,  # Блюда из птицы
            'fish': [],           # Пустой список
            'side_dishes': [],    # Пустой список
        }

        # Обновляем презентацию с найденными блюдами
        success = update_presentation_with_all_categories(template_path, all_dishes, output_path)
        
        if success:
            # Формируем подробное сообщение о результатах
            results = []
            if len(salads) > 0:
                results.append(f"Салаты и холодные закуски: {len(salads)} блюд")
            if len(first_courses) > 0:
                results.append(f"Первые блюда: {len(first_courses)} блюд")
            if len(meat_dishes) > 0:
                results.append(f"Блюда из мяса: {len(meat_dishes)} блюд")
            if len(poultry_dishes) > 0:
                results.append(f"Блюда из птицы: {len(poultry_dishes)} блюд")
            
            message = "Презентация создана!\n" + "\n".join(results)
            return True, message
        else:
            return False, "Ошибка при обновлении презентации"
            
    except Exception as e:
        return False, f"Ошибка: {str(e)}"
