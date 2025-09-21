#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Test script to verify the content of presentations created by the program
"""

from pathlib import Path
from pptx import Presentation
import sys

def analyze_presentation(pptx_path):
    """Analyze the contents of a PowerPoint presentation"""
    print(f"🔍 Анализируем презентацию: {pptx_path}")
    
    try:
        prs = Presentation(pptx_path)
        print(f"📋 Всего слайдов в презентации: {len(prs.slides)}")
        print()
        
        for i, slide in enumerate(prs.slides, 1):
            print(f"📄 Слайд {i}:")
            
            # Анализируем таблицы на слайде
            table_count = 0
            for shape in slide.shapes:
                if hasattr(shape, 'table'):
                    table_count += 1
                    table = shape.table
                    print(f"  📊 Таблица {table_count}: {len(table.rows)} строк, {len(table.columns)} столбцов")
                    
                    # Показываем первые несколько строк данных
                    for row_idx in range(min(5, len(table.rows))):
                        row_data = []
                        for cell in table.rows[row_idx].cells:
                            text = cell.text.strip()
                            if text:
                                row_data.append(text[:30] + ('...' if len(text) > 30 else ''))
                        if row_data:
                            print(f"    Строка {row_idx + 1}: {' | '.join(row_data)}")
                    
                    if len(table.rows) > 5:
                        print(f"    ... и еще {len(table.rows) - 5} строк")
                
                # Анализируем текстовые блоки
                if hasattr(shape, 'text'):
                    text = shape.text.strip()
                    if text and len(text) > 10:  # Игнорируем короткий текст
                        print(f"  📝 Текст: {text[:50]}...")
            
            if table_count == 0:
                print("  ❌ На слайде нет таблиц")
            
            print()
    
    except Exception as e:
        print(f"❌ Ошибка при анализе презентации: {e}")
        return False
    
    return True

def main():
    # Проверяем презентации на Рабочем столе
    desktop = Path.home() / "Desktop"
    
    presentations = list(desktop.glob("*.pptx"))
    
    if not presentations:
        print("❌ На Рабочем столе не найдено файлов презентаций (.pptx)")
        return
    
    print(f"🎯 Найдено презентаций: {len(presentations)}")
    print()
    
    for pptx_file in presentations:
        print("=" * 70)
        analyze_presentation(str(pptx_file))
        print("=" * 70)
        print()

if __name__ == "__main__":
    main()
