#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Диагностика проблем с таблицами на всех слайдах презентации
"""
import os
import sys
from pathlib import Path

sys.path.append(os.path.dirname(os.path.abspath(__file__)))

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    
    def analyze_all_slides():
        template_path = Path("templates/presentation_template.pptx")
        
        if not template_path.exists():
            print(f"❌ Шаблон не найден: {template_path}")
            return
        
        print("📊 АНАЛИЗ ВСЕХ СЛАЙДОВ И ТАБЛИЦ")
        print("=" * 60)
        
        prs = Presentation(str(template_path))
        print(f"Всего слайдов: {len(prs.slides)}")
        
        slide_names = ["Титульный", "Неизвестно", "Салаты", "Первые", "Мясо", "Птица", "Рыба", "Гарниры"]
        
        for slide_idx in range(len(prs.slides)):
            slide = prs.slides[slide_idx]
            slide_name = slide_names[slide_idx] if slide_idx < len(slide_names) else f"Слайд {slide_idx+1}"
            
            print(f"\n🎯 СЛАЙД {slide_idx+1} ({slide_name}):")
            
            tables = []
            for shape_idx, shape in enumerate(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table
                    table_info = {
                        'shape_idx': shape_idx,
                        'rows': len(table.rows),
                        'cols': len(table.columns),
                        'data_rows': len(table.rows) - 1
                    }
                    tables.append(table_info)
                    
                    print(f"   📋 Таблица {len(tables)} (Shape {shape_idx}):")
                    print(f"      Размер: {table_info['rows']} строк × {table_info['cols']} столбцов")
                    print(f"      Строк для данных: {table_info['data_rows']}")
                    
                    # Показываем заголовки первой строки
                    if table_info['rows'] > 0:
                        header_row = table.rows[0]
                        headers = []
                        for cell in header_row.cells:
                            headers.append(f"'{cell.text.strip()}'")
                        print(f"      Заголовки: {' | '.join(headers)}")
                    
                    # Показываем первые несколько строк данных
                    print(f"      Содержимое строк данных:")
                    for row_idx in range(1, min(4, table_info['rows'])):
                        row = table.rows[row_idx]
                        row_content = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_content.append(f"'{cell_text}'")
                            else:
                                row_content.append("'[пусто]'")
                        print(f"        Строка {row_idx}: {' | '.join(row_content)}")
            
            if not tables:
                print("   ❌ Таблиц не найдено")
            else:
                # Показываем какую таблицу выберет наша логика
                best_table = max(tables, key=lambda t: t['data_rows'])
                print(f"\n   🎯 ВЫБОР СИСТЕМЫ: Таблица {tables.index(best_table)+1}")
                print(f"      (Самая большая: {best_table['data_rows']} строк для данных)")
                
                if len(tables) > 1:
                    print(f"   ⚠️  ВНИМАНИЕ: На слайде {len(tables)} таблиц!")
                    for i, t in enumerate(tables, 1):
                        status = "👈 ВЫБРАНА" if t == best_table else ""
                        print(f"      Таблица {i}: {t['data_rows']} строк {status}")

    if __name__ == "__main__":
        analyze_all_slides()
        
except ImportError:
    print("❌ Модуль python-pptx не установлен")
except Exception as e:
    print(f"❌ Ошибка: {e}")
