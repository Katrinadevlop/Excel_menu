#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Диагностика проблем с извлечением и вставкой блюд
"""
import os
import sys
from pathlib import Path
import pandas as pd

sys.path.append(os.path.dirname(os.path.abspath(__file__)))

from presentation_handler import (
    extract_fish_dishes_from_column_e, 
    extract_side_dishes_from_excel,
    MenuItem
)

def analyze_fish_extraction_detailed():
    """Подробный анализ извлечения рыбных блюд"""
    print("🐟 ДЕТАЛЬНЫЙ АНАЛИЗ ИЗВЛЕЧЕНИЯ РЫБНЫХ БЛЮД")
    print("=" * 60)
    
    test_file = Path(r"C:\Users\katya\Downloads\Telegram Desktop\11 сентября - четверг.xlsx")
    
    if not test_file.exists():
        print(f"❌ Файл не найден: {test_file}")
        return
    
    # Анализируем сырые данные Excel
    try:
        xls = pd.ExcelFile(str(test_file))
        sheet_name = None
        for nm in xls.sheet_names:
            if 'касс' in str(nm).strip().lower():
                sheet_name = nm
                break
        if sheet_name is None and xls.sheet_names:
            sheet_name = xls.sheet_names[0]
            
        df = pd.read_excel(str(test_file), sheet_name=sheet_name, header=None, dtype=object)
        
        def row_text(row) -> str:
            parts = []
            for v in row:
                if pd.notna(v):
                    parts.append(str(v))
            return ' '.join(parts).strip()
        
        print(f"📄 Анализируем файл: {test_file.name}")
        print(f"📄 Лист: {sheet_name}")
        print(f"📊 Размер: {len(df)} строк, {len(df.columns)} столбцов")
        print()
        
        # Находим секцию рыбных блюд
        fish_start = None
        fish_end = None
        fish_columns = None
        
        for i in range(len(df)):
            row_content = row_text(df.iloc[i]).upper().replace('Ё', 'Е')
            
            if fish_start is None and ('БЛЮДА ИЗ РЫБЫ' in row_content or ('РЫБН' in row_content and 'БЛЮДА' in row_content)):
                fish_start = i
                print(f"🎯 Найден заголовок рыбы в строке {i+1}: {row_content}")
                
                # Определяем столбцы
                for col_idx in range(len(df.columns)):
                    if pd.notna(df.iloc[i, col_idx]):
                        cell_content = str(df.iloc[i, col_idx]).upper().replace('Ё', 'Е')
                        if 'БЛЮДА ИЗ РЫБЫ' in cell_content or ('РЫБН' in cell_content and 'БЛЮДА' in cell_content):
                            fish_columns = [col_idx, col_idx + 1, col_idx + 2] if col_idx + 2 < len(df.columns) else [col_idx]
                            print(f"📍 Столбцы рыбных блюд: {fish_columns}")
                            break
                continue
            
            if fish_start is not None and fish_end is None:
                if 'ГАРНИРЫ' in row_content or 'ГАРНИР' in row_content:
                    fish_end = i
                    print(f"🛑 Найден конец (гарниры) в строке {i+1}: {row_content}")
                    break
        
        if fish_start is None or fish_columns is None:
            print("❌ Не удалось найти секцию рыбных блюд")
            return
            
        if fish_end is None:
            fish_end = min(fish_start + 20, len(df))
            
        print(f"📍 Анализируем строки {fish_start+1} - {fish_end}")
        print()
        
        # Детально анализируем каждую строку в секции
        print("🔍 ДЕТАЛЬНЫЙ АНАЛИЗ СТРОК:")
        dishes_found = []
        
        for i in range(fish_start + 1, fish_end):
            if i >= len(df):
                break
                
            row = df.iloc[i]
            print(f"\n📋 СТРОКА {i+1}:")
            
            # Показываем содержимое всех ячеек
            all_cells = []
            for j, cell in enumerate(row):
                if pd.notna(cell):
                    cell_str = str(cell).strip()
                    if cell_str:
                        all_cells.append(f"Кол.{j}: '{cell_str}'")
                else:
                    all_cells.append(f"Кол.{j}: [пусто]")
            
            print(f"   Все ячейки: {' | '.join(all_cells)}")
            
            # Анализируем данные только из столбцов рыбных блюд
            fish_data = []
            for col_idx in fish_columns:
                if col_idx < len(df.columns) and pd.notna(df.iloc[i, col_idx]):
                    cell_text = str(df.iloc[i, col_idx]).strip()
                    if cell_text:
                        fish_data.append(cell_text)
            
            print(f"   Данные рыбы: {fish_data}")
            
            if fish_data:
                # Пытаемся определить название, вес, цену
                name = fish_data[0] if fish_data else ""
                weight = ""
                price = ""
                
                # Простая логика определения
                if name and not name.isupper() and len(name) > 2:
                    print(f"   ✅ Название: '{name}'")
                    
                    for value in fish_data[1:]:
                        if not weight and ('г' in value.lower() or 'мл' in value.lower() or 'шт' in value.lower()):
                            weight = value
                            print(f"   ⚖️  Вес: '{weight}'")
                        elif not price and any(char.isdigit() for char in value):
                            if not ('г' in value.lower() or 'мл' in value.lower() or 'шт' in value.lower()):
                                price = value
                                print(f"   💰 Цена: '{price}'")
                    
                    if name:
                        dishes_found.append(MenuItem(name=name, weight=weight, price=price))
                        print(f"   ✅ ДОБАВЛЕНО: {name} | {weight} | {price}")
                else:
                    print(f"   ❌ Пропущено (не подходит название): '{name}'")
            else:
                print(f"   ❌ Нет данных в столбцах рыбных блюд")
        
        print(f"\n🎯 ИТОГО НАЙДЕНО РЫБНЫХ БЛЮД: {len(dishes_found)}")
        for i, dish in enumerate(dishes_found, 1):
            print(f"   {i}. {dish.name} | {dish.weight} | {dish.price}")
            
    except Exception as e:
        print(f"❌ Ошибка анализа: {e}")

def analyze_garnish_extraction():
    """Анализ проблем с извлечением гарниров"""
    print("\n" + "=" * 60)
    print("🥔 АНАЛИЗ ПРОБЛЕМ С ГАРНИРАМИ")
    print("=" * 60)
    
    test_file = Path(r"C:\Users\katya\Downloads\Telegram Desktop\11 сентября - четверг.xlsx")
    
    if not test_file.exists():
        print(f"❌ Файл не найден: {test_file}")
        return
    
    print("🔍 Извлекаем гарниры...")
    garnishes = extract_side_dishes_from_excel(str(test_file))
    
    print(f"\n📊 Всего извлечено гарниров: {len(garnishes)}")
    print("\n🥔 СПИСОК ГАРНИРОВ:")
    
    problems = []
    for i, dish in enumerate(garnishes, 1):
        name_issue = not dish.name or len(dish.name) < 3
        weight_issue = not dish.weight
        price_issue = not dish.price
        
        status = "❌" if name_issue else "✅"
        print(f"   {i:2d}. {status} Название: '{dish.name or '[ОТСУТСТВУЕТ]'}'")
        print(f"       Вес: '{dish.weight or '[ОТСУТСТВУЕТ]'}'")
        print(f"       Цена: '{dish.price or '[ОТСУТСТВУЕТ]'}'")
        
        if name_issue:
            problems.append(f"Гарнир #{i}: отсутствует название")
        if weight_issue:
            problems.append(f"Гарнир #{i}: отсутствует вес")
        if price_issue:
            problems.append(f"Гарнир #{i}: отсутствует цена")
        print()
    
    if problems:
        print("⚠️  ОБНАРУЖЕННЫЕ ПРОБЛЕМЫ:")
        for problem in problems:
            print(f"   - {problem}")
    else:
        print("✅ Проблем с данными гарниров не обнаружено")

def test_table_capacity():
    """Тестирует ограничения таблиц в презентации"""
    print("\n" + "=" * 60)
    print("📋 АНАЛИЗ ОГРАНИЧЕНИЙ ТАБЛИЦ В ПРЕЗЕНТАЦИИ")
    print("=" * 60)
    
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        template_path = Path("templates/presentation_template.pptx")
        
        if not template_path.exists():
            print(f"❌ Шаблон не найден: {template_path}")
            return
        
        prs = Presentation(str(template_path))
        print(f"📊 Всего слайдов в шаблоне: {len(prs.slides)}")
        
        # Анализируем ключевые слайды
        slide_names = ["Неизвестно", "Неизвестно", "Салаты", "Первые", "Мясо", "Птица", "Рыба", "Гарниры"]
        
        for slide_idx in range(min(len(prs.slides), len(slide_names))):
            slide = prs.slides[slide_idx]
            slide_name = slide_names[slide_idx] if slide_idx < len(slide_names) else f"Слайд {slide_idx+1}"
            
            print(f"\n🎯 СЛАЙД {slide_idx+1} ({slide_name}):")
            
            tables = []
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table
                    tables.append({
                        'rows': len(table.rows),
                        'cols': len(table.columns),
                        'data_rows': len(table.rows) - 1  # минус заголовок
                    })
            
            if tables:
                for i, table_info in enumerate(tables, 1):
                    print(f"   📋 Таблица {i}: {table_info['rows']} строк, {table_info['cols']} столбцов")
                    print(f"      Строк для данных: {table_info['data_rows']}")
                    
                    if table_info['data_rows'] < 4:
                        print(f"      ⚠️  ОГРАНИЧЕНИЕ: поместится только {table_info['data_rows']} блюд!")
                    elif table_info['data_rows'] >= 10:
                        print(f"      ✅ ХОРОШО: поместится {table_info['data_rows']} блюд")
            else:
                print(f"   ❌ Таблиц не найдено")
                
    except ImportError:
        print("❌ Модуль python-pptx не установлен")
    except Exception as e:
        print(f"❌ Ошибка анализа презентации: {e}")

if __name__ == "__main__":
    analyze_fish_extraction_detailed()
    analyze_garnish_extraction()
    test_table_capacity()
