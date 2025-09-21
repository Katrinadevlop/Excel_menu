#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Тест полной презентации со всеми категориями блюд
"""
import os
import sys
from pathlib import Path

sys.path.append(os.path.dirname(os.path.abspath(__file__)))

from presentation_handler import create_presentation_with_excel_data

def test_full_presentation():
    """Тестирует создание полной презентации со всеми категориями"""
    print("🧪 ТЕСТ ПОЛНОЙ ПРЕЗЕНТАЦИИ СО ВСЕМИ КАТЕГОРИЯМИ")
    print("=" * 70)
    
    # Проблемный файл
    excel_path = Path(r"C:\Users\katya\Downloads\Telegram Desktop\4 сентября - четверг (2).xls")
    template_path = Path("templates/presentation_template.pptx")
    output_path = Path("test_full_presentation.pptx")
    
    if not excel_path.exists():
        print(f"❌ Excel файл не найден: {excel_path}")
        return
        
    if not template_path.exists():
        print(f"❌ Шаблон не найден: {template_path}")
        return
    
    print(f"📄 Excel файл: {excel_path.name}")
    print(f"📄 Шаблон: {template_path}")
    print(f"💾 Результат: {output_path}")
    
    # Создаем полную презентацию
    print(f"\n🎯 СОЗДАНИЕ ПОЛНОЙ ПРЕЗЕНТАЦИИ")
    try:
        success, message = create_presentation_with_excel_data(
            str(template_path),
            str(excel_path),
            str(output_path)
        )
        
        if success:
            print(f"✅ Презентация создана: {output_path}")
            print(f"📝 Сообщение: {message}")
            
            if output_path.exists():
                size = output_path.stat().st_size
                print(f"📏 Размер файла: {size:,} байт")
                
                # Анализируем созданную презентацию
                analyze_full_presentation(str(output_path))
        else:
            print(f"❌ Ошибка создания презентации: {message}")
            
    except Exception as e:
        print(f"❌ Исключение: {e}")
        import traceback
        traceback.print_exc()

def analyze_full_presentation(pptx_path: str):
    """Анализирует все слайды полной презентации"""
    print(f"\n🔍 АНАЛИЗ ПОЛНОЙ ПРЕЗЕНТАЦИИ")
    print("=" * 70)
    
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        prs = Presentation(pptx_path)
        print(f"📊 Общее количество слайдов: {len(prs.slides)}")
        
        slide_categories = [
            "Титульный слайд",
            "Салаты и холодные закуски", 
            "Первые блюда", 
            "Мясные блюда",
            "Блюда из птицы",
            "Рыбные блюда", 
            "Гарниры"
        ]
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            category = slide_categories[slide_idx - 1] if slide_idx <= len(slide_categories) else f"Слайд {slide_idx}"
            print(f"\n🎯 СЛАЙД {slide_idx} ({category}):")
            
            tables_found = 0
            dishes_found = 0
            
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    tables_found += 1
                    table = shape.table
                    print(f"📋 Таблица {tables_found}: {len(table.rows)} строк × {len(table.columns)} столбцов")
                    
                    # Подсчитываем строки с данными (пропускаем заголовок)
                    for i in range(1, min(6, len(table.rows))):  # Показываем первые 5 строк данных
                        row = table.rows[i]
                        if len(row.cells) >= 3:
                            name = row.cells[0].text.strip()
                            weight = row.cells[1].text.strip()
                            price = row.cells[2].text.strip()
                            
                            if name:  # Если есть название блюда
                                dishes_found += 1
                                print(f"   {i}. '{name}' | '{weight}' | '{price}'")
                            else:
                                print(f"   {i}. [пусто] | '{weight}' | '{price}'")
                    
                    if len(table.rows) > 6:
                        print(f"   ... (еще {len(table.rows) - 6} строк)")
            
            if not tables_found:
                print("   📋 Таблиц не найдено")
            else:
                if dishes_found > 0:
                    print(f"   ✅ Найдено блюд с названиями: {dishes_found}")
                else:
                    print(f"   ❌ Блюда с названиями НЕ найдены")
                
    except ImportError:
        print("❌ Модуль python-pptx не установлен")
    except Exception as e:
        print(f"❌ Ошибка анализа презентации: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    test_full_presentation()
