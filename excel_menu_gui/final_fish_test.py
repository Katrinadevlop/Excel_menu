#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Финальный тест извлечения рыбных блюд и создания презентации
"""
import os
import sys
from pathlib import Path

# Добавляем текущую папку в путь для импорта наших модулей
sys.path.append(os.path.dirname(os.path.abspath(__file__)))

from presentation_handler import (
    extract_fish_dishes_from_column_e, 
    create_presentation_with_fish_and_side_dishes,
    MenuItem
)

def run_complete_test():
    """
    Полный тест: извлечение рыбных блюд + создание презентации + анализ результата
    """
    print("🧪 ПОЛНЫЙ ТЕСТ ФУНКЦИОНАЛЬНОСТИ РЫБНЫХ БЛЮД")
    print("=" * 70)
    
    # 1. Найдем Excel файл с рыбными блюдами
    downloads_path = Path(r"C:\Users\katya\Downloads\Telegram Desktop")
    excel_files = list(downloads_path.glob('*.xlsx')) + list(downloads_path.glob('*.xls'))
    
    test_excel = None
    best_fish_count = 0
    
    print("🔍 Поиск подходящего Excel файла...")
    for excel_file in excel_files[:5]:  # Проверяем первые 5 файлов
        try:
            dishes = extract_fish_dishes_from_column_e(str(excel_file))
            if len(dishes) > best_fish_count:
                best_fish_count = len(dishes)
                test_excel = excel_file
        except:
            continue
    
    if not test_excel:
        print("❌ Не найден подходящий Excel файл")
        return
    
    print(f"✅ Выбран файл: {test_excel.name}")
    print(f"   Рыбных блюд в файле: {best_fish_count}")
    
    # 2. Извлекаем рыбные блюда подробно
    print(f"\n📊 ПОДРОБНЫЙ АНАЛИЗ РЫБНЫХ БЛЮД")
    print("-" * 50)
    
    fish_dishes = extract_fish_dishes_from_column_e(str(test_excel))
    
    if not fish_dishes:
        print("❌ Рыбные блюда не извлечены")
        return
    
    print(f"✅ Успешно извлечено: {len(fish_dishes)} блюд")
    
    # Показываем каждое блюдо с деталями
    for i, dish in enumerate(fish_dishes, 1):
        print(f"\n🐟 БЛЮДО {i}:")
        print(f"   Название: '{dish.name}'")
        print(f"   Вес:      '{dish.weight or 'не указан'}'")
        print(f"   Цена:     '{dish.price or 'не указана'}'")
        
        # Проверяем качество данных
        issues = []
        if not dish.name or len(dish.name) < 3:
            issues.append("короткое/отсутствует название")
        if not dish.weight:
            issues.append("отсутствует вес")
        if not dish.price:
            issues.append("отсутствует цена")
            
        if issues:
            print(f"   ⚠️  Проблемы: {', '.join(issues)}")
        else:
            print(f"   ✅ Данные полные")
    
    # 3. Создаем презентацию
    print(f"\n🎯 СОЗДАНИЕ ПРЕЗЕНТАЦИИ")
    print("-" * 50)
    
    template_path = Path("templates/presentation_template.pptx")
    output_path = Path("test_final_fish_presentation.pptx")
    
    if not template_path.exists():
        print(f"❌ Шаблон не найден: {template_path}")
        return
    
    print(f"📄 Используем шаблон: {template_path}")
    print(f"💾 Создаем презентацию: {output_path}")
    
    try:
        success, message = create_presentation_with_fish_and_side_dishes(
            str(template_path),
            str(test_excel),
            str(output_path)
        )
        
        if success:
            print(f"✅ Презентация создана успешно!")
            print(f"📝 Сообщение: {message}")
            
            if output_path.exists():
                size = output_path.stat().st_size
                print(f"📏 Размер файла: {size:,} байт")
                
                # Анализируем что попало в презентацию
                analyze_presentation_content(str(output_path), fish_dishes)
            else:
                print("❌ Файл презентации не создался")
        else:
            print(f"❌ Ошибка создания презентации: {message}")
            
    except Exception as e:
        print(f"❌ Исключение: {e}")
        import traceback
        traceback.print_exc()

def analyze_presentation_content(pptx_path: str, original_dishes: list):
    """
    Анализирует содержимое созданной презентации
    """
    print(f"\n🔍 АНАЛИЗ СОЗДАННОЙ ПРЕЗЕНТАЦИИ")
    print("-" * 50)
    
    try:
        from pptx import Presentation
        
        prs = Presentation(pptx_path)
        print(f"📊 Слайдов в презентации: {len(prs.slides)}")
        
        if len(prs.slides) >= 6:
            slide_6 = prs.slides[5]  # 6-й слайд (индекс 5)
            print(f"🎯 Анализируем 6-й слайд (где должны быть рыбные блюда):")
            
            # Ищем все текстовые элементы на слайде
            text_content = []
            table_found = False
            
            for shape in slide_6.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())
                
                # Проверяем есть ли таблицы
                if hasattr(shape, "table"):
                    table_found = True
                    table = shape.table
                    print(f"📋 Найдена таблица: {len(table.rows)} строк, {len(table.columns)} столбцов")
                    
                    # Показываем содержимое таблицы
                    for i, row in enumerate(table.rows):
                        row_content = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_content.append(cell_text)
                        if row_content:
                            print(f"   Строка {i+1}: {' | '.join(row_content)}")
            
            if text_content:
                print(f"📝 Текст на слайде:")
                for text in text_content:
                    if len(text) > 100:
                        print(f"   {text[:100]}...")
                    else:
                        print(f"   {text}")
            
            # Проверяем, попали ли наши рыбные блюда в презентацию
            slide_text = ' '.join(text_content).upper()
            dishes_found = 0
            
            print(f"\n🎯 ПРОВЕРКА ПОПАДАНИЯ БЛЮД В ПРЕЗЕНТАЦИЮ:")
            for dish in original_dishes:
                dish_name_upper = dish.name.upper()
                if dish_name_upper in slide_text:
                    dishes_found += 1
                    print(f"   ✅ '{dish.name}' - найдено в презентации")
                else:
                    print(f"   ❌ '{dish.name}' - НЕ найдено в презентации")
            
            print(f"\n📈 ИТОГО: {dishes_found}/{len(original_dishes)} блюд попало в презентацию")
            
            if dishes_found == len(original_dishes):
                print("🎉 ВСЕ РЫБНЫЕ БЛЮДА УСПЕШНО ПЕРЕНЕСЕНЫ В ПРЕЗЕНТАЦИЮ!")
            elif dishes_found > 0:
                print("⚠️  Часть блюд перенесена, но не все")
            else:
                print("❌ Рыбные блюда не найдены в презентации")
                
        else:
            print(f"❌ В презентации недостаточно слайдов (нужен 6-й слайд)")
            
    except ImportError:
        print("❌ Модуль python-pptx не установлен, анализ презентации невозможен")
    except Exception as e:
        print(f"❌ Ошибка анализа презентации: {e}")

if __name__ == "__main__":
    run_complete_test()
